"""
TIC Copilot — Claude x GPT (IAedu)
==================================================================
Assistente multi-modelo para o Technical Innovation Center (CAD,
3D Printing, 3D Reality, Mendix). Roteia cada pedido para o modelo
mais forte (Claude Opus 4.7 vs GPT-5.5), com modos Duelo e
Colaboracao. Inclui:

  - Login com password (PBKDF2, stdlib)
  - HISTORICO de conversas persistente por utilizador (criar /
    retomar / apagar / exportar para Markdown)
  - Memoria "master" por utilizador (resumo persistente do perfil)
  - Upload de PDF / PPTX / imagem
  - OCR (tesseract) + VLM gratuito (Hugging Face router) para
    interpretar imagens (diagramas, pecas, screenshots Mendix...)
  - RAG HIBRIDO estado-da-arte: BM25 + embeddings densos (FAISS)
    fundidos por Reciprocal Rank Fusion + reranker neural opcional
  - PESQUISA WEB (DuckDuckGo, sem chave) com leitura das paginas
  - ENTRADA POR VOZ: gravacao no browser + speech-to-text local
    (faster-whisper, gratuito)
  - Acoes rapidas TIC, dashboard de custos da sessao

Correr:
  pip install -r requirements.txt
  # OCR precisa do binario do sistema:
  #   sudo apt-get install tesseract-ocr tesseract-ocr-por
  # VLM (opcional): define HF_TOKEN (token gratuito de huggingface.co)
  streamlit run app.py
"""

from __future__ import annotations

import base64
import datetime as dt
import hashlib
import hmac
import io
import json
import os
import queue
import re
import secrets
import threading
import time
from pathlib import Path

import requests
import streamlit as st
import streamlit.components.v1 as components
st.set_page_config(page_icon="Gemini_Generated_Image_o58khxo58khxo58k.png")
# ---------------------------------------------------------------------------
# Dependencias opcionais (extracao/OCR) - degradacao graciosa se faltarem
# ---------------------------------------------------------------------------
try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except Exception:
    HAS_FITZ = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

try:
    from PIL import Image
    HAS_PIL = True
except Exception:
    HAS_PIL = False

OCR_LANG = "eng"
try:
    import pytesseract
    pytesseract.get_tesseract_version()
    langs = set(pytesseract.get_languages(config=""))
    OCR_LANG = "+".join([l for l in ("por", "eng") if l in langs]) or "eng"
    HAS_OCR = HAS_PIL
except Exception:
    HAS_OCR = False

# ---------------------------------------------------------------------------
# Dependencias opcionais para RAG hibrido
# ---------------------------------------------------------------------------
try:
    import numpy as np
    HAS_NUMPY = True
except Exception:
    np = None
    HAS_NUMPY = False

try:
    import faiss
    HAS_FAISS = True
except Exception:
    faiss = None
    HAS_FAISS = False

try:
    from sentence_transformers import SentenceTransformer, CrossEncoder
    HAS_SENTENCE_TRANSFORMERS = True
except Exception:
    SentenceTransformer = None
    CrossEncoder = None
    HAS_SENTENCE_TRANSFORMERS = False

try:
    from rank_bm25 import BM25Okapi
    HAS_BM25 = True
except Exception:
    BM25Okapi = None
    HAS_BM25 = False

# ---------------------------------------------------------------------------
# Dependencias opcionais: pesquisa web (DuckDuckGo, sem chave de API)
# ---------------------------------------------------------------------------
HAS_DDG = False
try:
    from ddgs import DDGS
    HAS_DDG = True
except Exception:
    try:
        from duckduckgo_search import DDGS
        HAS_DDG = True
    except Exception:
        DDGS = None

# ---------------------------------------------------------------------------
# Dependencias opcionais: speech-to-text local (faster-whisper)
# ---------------------------------------------------------------------------
try:
    from faster_whisper import WhisperModel
    HAS_WHISPER = True
except Exception:
    WhisperModel = None
    HAS_WHISPER = False

APP_DIR = Path(__file__).resolve().parent
DATA_DIR = APP_DIR / "data"
USERS_FILE = APP_DIR / "users.json"
USERNAME_RE = re.compile(r"^[a-z0-9_.-]{3,32}$")

PER_FILE_CAP = 20000
TOTAL_CAP = 40000
SESSION_TRANSCRIPT_CAP = 24000  # janela maxima da conversa atual enviada aos modelos
MEMORY_UPDATE_EVERY = 3   # turnos entre atualizacoes automaticas da memoria
MAX_CONVERSATIONS_LISTED = 18

# RAG hibrido para anexos grandes/persistentes da sessao
ATTACHMENT_CHUNK_SIZE = 1800
ATTACHMENT_CHUNK_OVERLAP = 250
ATTACHMENT_CONTEXT_CAP = 36000
ATTACHMENT_TOP_K_DEFAULT = 6
ATTACHMENT_CANDIDATE_MULTIPLIER = 6
RRF_K = 60  # constante classica do Reciprocal Rank Fusion
EMBEDDING_MODEL_OPTIONS = [
    "BAAI/bge-m3",
    "intfloat/multilingual-e5-large-instruct",
    "sentence-transformers/paraphrase-multilingual-mpnet-base-v2",
]
RERANKER_MODEL_OPTIONS = [
    "BAAI/bge-reranker-v2-m3",
    "cross-encoder/ms-marco-MiniLM-L-6-v2",
]

# Pesquisa web
WEB_MAX_RESULTS = 6
WEB_FETCH_PAGES = 2
WEB_PAGE_CAP = 4500
WEB_CONTEXT_CAP = 16000

# VLM gratuito (Hugging Face Inference Router - token gratuito chega)
HF_VLM_MODEL = "Qwen/Qwen2.5-VL-7B-Instruct"
HF_ROUTER_URL = "https://router.huggingface.co/v1/chat/completions"

# Speech-to-text
WHISPER_SIZES = ["tiny", "base", "small"]
WHISPER_DEFAULT = "base"

# ---------------------------------------------------------------------------
# Configuracao dos dois agentes (IAedu)
# ---------------------------------------------------------------------------
MODELS = {
    "claude": {
        "label": "Claude Opus 4.7", "short": "Opus 4.7", "color": "#D2A24C",
        "endpoint": "https://api.iaedu.pt/agent-chat//api/v1/agent/cmoss7l0f658oko01vk2egfpg/stream",
        "channel_id": "cmpws0w210p11i601iwo0bxpe",
        "env": "IAEDU_CLAUDE_KEY", "price_in": 5.0, "price_out": 25.0,
    },
    "gpt": {
        "label": "GPT-5.5", "short": "GPT-5.5", "color": "#46B98C",
        "endpoint": "https://api.iaedu.pt/agent-chat//api/v1/agent/cmor5objoex9gfp01vm7p95jh/stream",
        "channel_id": "cmpzy2p8j05wanr010emnsh2w",
        "env": "IAEDU_GPT_KEY", "price_in": 5.0, "price_out": 30.0,
    },
}

# ---------------------------------------------------------------------------
# Regras de roteamento (benchmarks Opus 4.7 vs GPT-5.5 + dominio TIC)
# ---------------------------------------------------------------------------
ROUTING_RULES = [
    {"id": "terminal", "model": "gpt", "label": "Agentes de terminal / CLI / DevOps",
     "bench": "Terminal-Bench 2.0 - 82.7% vs 69.4%",
     "keywords": ["terminal", "shell", "bash", "cli", "comando", "linha de comando", "devops",
                  "infra", "infraestrutura", "docker", "kubernetes", "ssh", "ci/cd",
                  "pipeline ci", "cron", "ansible", "makefile"]},
    {"id": "web", "model": "gpt", "label": "Pesquisa web / browsing / info atual",
     "bench": "BrowseComp - 90.1% (Pro) vs 79.3%",
     "keywords": ["pesquisa web", "pesquisar na web", "pesquisa na internet", "browse", "navegar",
                  "internet", "noticias", "notícias", "atual", "ultimo", "último", "recente",
                  "cotacao", "cotação", "preco atual", "preço atual", "search the web", "google",
                  "hoje em dia", "lancamento", "lançamento", "novidades"]},
    {"id": "longcontext", "model": "gpt", "label": "Contexto longo (>128K tokens)",
     "bench": "MRCR v2 512K-1M - 74.0% vs 32.2%",
     "keywords": ["contexto longo", "documento inteiro", "monorepo", "corpus", "ficheiro gigante",
                  "long context", "1m tokens", "ano de registos", "logs inteiros",
                  "analisa este documento longo", "milhares de paginas"]},
    {"id": "math", "model": "gpt", "label": "Matematica de fronteira / competicao",
     "bench": "FrontierMath Tier 4 - 35.4% vs 22.9%",
     "keywords": ["prova matematica", "prova matemática", "demonstra que", "teorema", "olimpiada",
                  "olimpíada", "competicao matematica", "integral", "equacao diferencial",
                  "equação diferencial", "algebra abstrata", "prove that", "calculo avancado", "lema"]},
    {"id": "abstract", "model": "gpt", "label": "Raciocinio abstrato / padroes",
     "bench": "ARC-AGI-2 - 85.0% vs 75.8%",
     "keywords": ["padrao", "padrão", "puzzle", "sequencia logica", "sequência lógica",
                  "raciocinio abstrato", "arc-agi", "quebra-cabeca", "adivinha a regra",
                  "qual o proximo", "qual o próximo"]},
    {"id": "office", "model": "gpt", "label": "Trabalho de escritorio / documentos",
     "bench": "OfficeQA Pro - 54.1% vs 43.6%  ·  GDPval 84.9% vs 80.3%",
     "keywords": ["powerpoint", "slides", "excel", "folha de calculo", "folha de cálculo",
                  "relatorio de escritorio", "memorando", "ata", "documento word"]},
    {"id": "mendix", "model": "claude", "label": "Mendix / low-code / arquitetura de apps",
     "bench": "SWE-Bench Pro - forte em codigo e arquitetura",
     "keywords": ["mendix", "low-code", "lowcode", "low code", "microflow", "nanoflow",
                  "domain model", "modelo de dominio", "entidade mendix", "modulo mendix",
                  "workflow", "xpath", "ocl", "arquitetura da aplicacao",
                  "arquitetura da aplicação", "logica de negocio", "lógica de negócio",
                  "alocacao automatica", "alocação automática", "regras de negocio"]},
    {"id": "cad3d", "model": "claude", "label": "CAD / Impressao 3D / 3D Reality",
     "bench": "Raciocinio tecnico profundo de engenharia",
     "keywords": ["cad", "catia", "solidworks", "autocad", "nx", "creo", "impressao 3d",
                  "impressão 3d", "impressora 3d", "stl", "step", "g-code", "gcode", "fdm",
                  "sla", "sls", "slicer", "tolerancia", "tolerância", "peca", "peça",
                  "prototipo", "protótipo", "realidade aumentada", "realidade virtual",
                  "ar/vr", "digital twin", "gemeo digital", "gémeo digital", "scan 3d"]},
    {"id": "leadership", "model": "claude", "label": "Lideranca / gestao de equipas",
     "bench": "Escrita e raciocinio organizacional",
     "keywords": ["lideranca", "liderança", "equipa", "gestao de equipa", "gestão de equipa",
                  "chefiar", "reuniao", "reunião", "one-on-one", "feedback", "onboarding",
                  "organizacao do departamento", "organização do departamento", "kpi", "kpis",
                  "okr", "delegar", "motivar", "avaliacao de desempenho",
                  "avaliação de desempenho"]},
    {"id": "pr_code", "model": "claude", "label": "Resolucao de PRs / refactor / codebase",
     "bench": "SWE-Bench Pro - 64.3% vs 58.6%",
     "keywords": ["bug", "corrige", "fix", "refactor", "refatora", "pull request", "code review",
                  "reve o codigo", "revê o código", "codebase", "repositorio", "repositório",
                  "stack trace", "erro no codigo", "erro no código", "debugar", "depurar",
                  "migracao de codigo", "traceback", "exception"]},
    {"id": "mcp", "model": "claude", "label": "Orquestracao de ferramentas / MCP / agentes",
     "bench": "MCP Atlas - 79.1% vs 75.3%",
     "keywords": ["mcp", "orquestra", "tool use", "agente multi-step", "chama a api", "encadeia",
                  "workflow de agentes", "function calling", "multiplas ferramentas",
                  "múltiplas ferramentas", "fastapi"]},
    {"id": "finance", "model": "claude", "label": "Analise financeira",
     "bench": "FinanceAgent v1.1 - 64.4% vs 60.0%",
     "keywords": ["financeiro", "financas", "finanças", "investimento", "balanco", "balanço",
                  "dcf", "valuation", "fluxo de caixa", "acoes", "ações", "portfolio", "portfólio",
                  "analise financeira", "demonstracoes financeiras", "orcamento", "orçamento",
                  "custo por peca", "custo por peça"]},
    {"id": "academic", "model": "claude", "label": "Raciocinio academico profundo",
     "bench": "Humanity's Last Exam - 54.7% vs 52.2% (c/ ferramentas)",
     "keywords": ["tese", "paper", "artigo cientifico", "artigo científico", "investigacao",
                  "investigação", "academico", "académico", "prova conceptual",
                  "explica em profundidade", "raciocinio profundo", "dissertacao", "dissertação",
                  "estado da arte", "review critica"]},
]

DEFAULT_MODEL = "claude"
DEFAULT_REASON = ("Sem sinal forte de categoria &mdash; Claude por defeito "
                  "(saida ~17% mais barata e forte em raciocinio geral / codigo).")
LONG_INPUT_CHARS = 6000

# Acoes rapidas para o contexto TIC (Autoeuropa)
QUICK_ACTIONS = [
    ("Mendix: alocacao automatica",
     "Estou a desenvolver em Mendix a logica de alocacao automatica de tecnico responsavel e "
     "impressora 3D para pedidos de impressao de pecas, otimizando o tempo de espera. Propoe a "
     "arquitetura (domain model, microflows, filas de prioridade) e a logica de negocio passo a passo."),
    ("Otimizar fila de impressao 3D",
     "Como devo modelar e otimizar a fila de pedidos de impressao 3D do departamento (varias "
     "impressoras, materiais e prioridades diferentes) para minimizar tempos de espera? Inclui "
     "algoritmos praticos e como implementa-los em Mendix."),
    ("Estruturar o departamento TIC",
     "Vou chefiar um novo departamento TIC que agrega CAD, 3D Printing e 3D Reality. Ajuda-me a "
     "estruturar a equipa, definir responsabilidades, KPIs e rituais de gestao (reunioes, "
     "prioritizacao de pedidos)."),
    ("Rever requisitos da app",
     "Faz uma revisao critica dos requisitos da minha aplicacao de gestao de pedidos CAD/3D "
     "(alocacao automatica de tecnico e impressora, otimizacao do tempo de espera): lacunas, "
     "casos extremos e melhorias."),
]


# ---------------------------------------------------------------------------
# Roteador
# ---------------------------------------------------------------------------
def route(message: str, total_len: int | None = None) -> dict:
    text = message.lower()
    scores: dict[str, int] = {}
    for rule in ROUTING_RULES:
        hits = sum(1 for kw in rule["keywords"] if kw in text)
        if hits:
            scores[rule["id"]] = hits
    if (total_len or len(message)) > LONG_INPUT_CHARS:
        scores["longcontext"] = scores.get("longcontext", 0) + 2

    if not scores:
        return {"model": DEFAULT_MODEL, "rule_id": None, "label": "Geral",
                "bench": "-", "reason": DEFAULT_REASON, "scores": {}}

    best_id, best = None, -1
    for rule in ROUTING_RULES:
        s = scores.get(rule["id"], 0)
        if s > best:
            best, best_id = s, rule["id"]

    rule = next(r for r in ROUTING_RULES if r["id"] == best_id)
    reason = (f"Categoria detetada: <b>{rule['label']}</b> &rarr; "
              f"{MODELS[rule['model']]['label']} ({rule['bench']})")
    return {"model": rule["model"], "rule_id": rule["id"], "label": rule["label"],
            "bench": rule["bench"], "reason": reason, "scores": scores}


# ---------------------------------------------------------------------------
# Autenticacao (PBKDF2 + users.json)
# ---------------------------------------------------------------------------
def _load_users() -> dict:
    if USERS_FILE.exists():
        try:
            return json.loads(USERS_FILE.read_text("utf-8"))
        except Exception:
            return {}
    return {}


def _save_users(users: dict):
    USERS_FILE.write_text(json.dumps(users, indent=2, ensure_ascii=False), "utf-8")


def _hash_pw(pw: str, salt: str) -> str:
    return hashlib.pbkdf2_hmac("sha256", pw.encode(), bytes.fromhex(salt), 200_000).hex()


def create_user(username: str, name: str, password: str):
    users = _load_users()
    salt = secrets.token_hex(16)
    users[username] = {"salt": salt, "hash": _hash_pw(password, salt),
                       "name": name, "created": dt.datetime.now().isoformat(timespec="seconds")}
    _save_users(users)


def verify_user(username: str, password: str) -> bool:
    u = _load_users().get(username)
    if not u:
        return False
    return hmac.compare_digest(u["hash"], _hash_pw(password, u["salt"]))


# ---------------------------------------------------------------------------
# Memoria do utilizador
# ---------------------------------------------------------------------------
def _user_dir(username: str) -> Path:
    d = DATA_DIR / username
    d.mkdir(parents=True, exist_ok=True)
    return d


def load_memory(username: str) -> dict:
    f = _user_dir(username) / "memory.json"
    if f.exists():
        try:
            return json.loads(f.read_text("utf-8"))
        except Exception:
            pass
    return {"master": "", "updated_at": None, "turns_since": 0}


def save_memory(username: str, mem: dict):
    (_user_dir(username) / "memory.json").write_text(
        json.dumps(mem, indent=2, ensure_ascii=False), "utf-8")


def append_log(username: str, prompt: str, mode: str):
    line = {"ts": dt.datetime.now().isoformat(timespec="seconds"),
            "mode": mode, "prompt": prompt[:1500]}
    with open(_user_dir(username) / "log.jsonl", "a", encoding="utf-8") as fh:
        fh.write(json.dumps(line, ensure_ascii=False) + "\n")


def recent_log(username: str, n: int = 12) -> list[str]:
    f = _user_dir(username) / "log.jsonl"
    if not f.exists():
        return []
    lines = f.read_text("utf-8").splitlines()[-n:]
    out = []
    for ln in lines:
        try:
            out.append(json.loads(ln)["prompt"])
        except Exception:
            pass
    return out


def summarize_memory(username: str):
    """Atualiza o resumo master usando o Claude (best-effort)."""
    mem = load_memory(username)
    recents = recent_log(username, 14)
    if not recents:
        return
    prompt = (
        "Es um sistema de memoria de longo prazo. Atualiza um resumo CONCISO (PT-PT, "
        "no maximo 160 palavras) sobre o utilizador: quem parece ser, areas e temas "
        "recorrentes, e o tipo de pedidos que costuma fazer. Combina o RESUMO ATUAL com "
        "as INTERACOES RECENTES. Nao inventes nada; se nao houver dados, mantem vago. "
        "Devolve apenas o resumo, sem preambulos.\n\n"
        f"=== RESUMO ATUAL ===\n{mem.get('master') or '(vazio)'}\n\n"
        f"=== INTERACOES RECENTES ===\n- " + "\n- ".join(recents)
    )
    try:
        text = "".join(respond("claude", prompt, new_thread_id(), api_key_for("claude")))
        text = text.strip()
        if text and not text.lower().startswith(("> erro", "> falha", "> limite")):
            mem["master"] = text
            mem["updated_at"] = dt.datetime.now().isoformat(timespec="seconds")
            mem["turns_since"] = 0
            save_memory(username, mem)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Historico de conversas (persistente por utilizador)
# ---------------------------------------------------------------------------
def _conv_dir(username: str) -> Path:
    d = _user_dir(username) / "conversations"
    d.mkdir(parents=True, exist_ok=True)
    return d


def list_conversations(username: str) -> list[dict]:
    items = []
    for f in _conv_dir(username).glob("*.json"):
        try:
            d = json.loads(f.read_text("utf-8"))
            items.append({"id": d.get("id", f.stem),
                          "title": d.get("title", "(sem titulo)"),
                          "updated": d.get("updated", ""),
                          "n": len(d.get("messages", []))})
        except Exception:
            continue
    items.sort(key=lambda x: x.get("updated", ""), reverse=True)
    return items


def save_current_conversation(username: str):
    """Grava a conversa atual em disco (chamado apos cada turno)."""
    msgs = st.session_state.get("messages", [])
    if not msgs:
        return
    cid = st.session_state.get("conv_id")
    if not cid:
        cid = dt.datetime.now().strftime("%Y%m%d-%H%M%S") + "-" + secrets.token_hex(3)
        st.session_state.conv_id = cid
    title = st.session_state.get("conv_title")
    if not title:
        first = next((m.get("content", "") for m in msgs if m.get("role") == "user"), "Conversa")
        first = " ".join(first.split())
        title = (first[:64] + "...") if len(first) > 64 else (first or "Conversa")
        st.session_state.conv_title = title
    created = st.session_state.get("conv_created") or dt.datetime.now().isoformat(timespec="seconds")
    st.session_state.conv_created = created
    data = {
        "id": cid, "title": title, "created": created,
        "updated": dt.datetime.now().isoformat(timespec="seconds"),
        "messages": msgs,
        "threads": st.session_state.get("threads", {}),
        "rag_docs": [{"name": d.get("name"), "text": d.get("text"), "hash": d.get("hash")}
                     for d in st.session_state.get("rag_docs", [])],
    }
    try:
        (_conv_dir(username) / f"{cid}.json").write_text(
            json.dumps(data, ensure_ascii=False), "utf-8")
    except Exception:
        pass


def load_conversation(username: str, cid: str) -> bool:
    f = _conv_dir(username) / f"{cid}.json"
    if not f.exists():
        return False
    try:
        d = json.loads(f.read_text("utf-8"))
    except Exception:
        return False
    st.session_state.messages = d.get("messages", [])
    threads = d.get("threads", {}) or {}
    st.session_state.threads = {k: threads.get(k) or new_thread_id() for k in MODELS}
    st.session_state.conv_id = d.get("id", cid)
    st.session_state.conv_title = d.get("title", "")
    st.session_state.conv_created = d.get("created")
    st.session_state.rag_docs = d.get("rag_docs", []) or []
    st.session_state.rag_doc_hashes = {x.get("hash") for x in st.session_state.rag_docs if x.get("hash")}
    st.session_state.uploader_key = st.session_state.get("uploader_key", 0) + 1
    return True


def delete_conversation(username: str, cid: str):
    f = _conv_dir(username) / f"{cid}.json"
    try:
        f.unlink(missing_ok=True)
    except Exception:
        pass


def conversation_markdown() -> str:
    """Exporta a conversa atual para Markdown."""
    lines = [f"# {st.session_state.get('conv_title') or 'Conversa'}",
             f"_Exportado em {dt.datetime.now().strftime('%Y-%m-%d %H:%M')}_", ""]
    for msg in st.session_state.get("messages", []):
        if msg.get("role") == "user":
            lines.append("## Pedido")
            lines.append(msg.get("content", ""))
            if msg.get("attachments"):
                lines.append("*Anexos: " + ", ".join(msg["attachments"]) + "*")
        else:
            mode = msg.get("mode")
            if mode == "single":
                lines.append(f"## Resposta ({MODELS.get(msg.get('model', ''), {}).get('label', '?')})")
                lines.append(msg.get("content", ""))
            elif mode == "duel":
                lines.append("## Duelo")
                lines.append("### Claude\n" + msg.get("data", {}).get("claude", ""))
                lines.append("### GPT\n" + msg.get("data", {}).get("gpt", ""))
            elif mode == "collab":
                lines.append("## Colaboracao")
                lines.append("### Rascunho\n" + msg.get("draft", ""))
                lines.append("### Revisao\n" + msg.get("critique", ""))
                if msg.get("synthesis"):
                    lines.append("### Versao final\n" + msg.get("synthesis", ""))
        lines.append("")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Extracao de ficheiros (PDF / PPTX / imagem) + OCR + VLM
# ---------------------------------------------------------------------------
def hf_token() -> str | None:
    """Token gratuito do Hugging Face (HF_TOKEN) para o VLM."""
    tok = os.getenv("HF_TOKEN")
    if tok:
        return tok.strip()
    try:
        if "HF_TOKEN" in st.secrets:
            return str(st.secrets["HF_TOKEN"]).strip()
    except Exception:
        pass
    return None


def describe_image_vlm(data: bytes, prompt: str = "", mime: str = "image/png") -> str | None:
    """Interpreta a imagem com um VLM gratuito (Hugging Face router).

    Devolve None se nao houver token ou se a chamada falhar — o chamador
    deve degradar para OCR.
    """
    token = hf_token()
    if not token:
        return None
    instr = ("Descreve esta imagem em PT-PT de forma tecnica e util: o que mostra, "
             "texto visivel, diagramas, medidas, e qualquer detalhe relevante para "
             "engenharia (CAD, impressao 3D, Mendix, esquemas).")
    if prompt:
        instr += f"\nContexto do pedido do utilizador: {prompt[:500]}"
    b64 = base64.b64encode(data).decode()
    payload = {
        "model": HF_VLM_MODEL,
        "max_tokens": 700,
        "messages": [{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                {"type": "text", "text": instr},
            ],
        }],
    }
    try:
        r = requests.post(HF_ROUTER_URL, json=payload, timeout=90,
                          headers={"Authorization": f"Bearer {token}"})
        r.raise_for_status()
        out = r.json()["choices"][0]["message"]["content"]
        return (out or "").strip() or None
    except Exception:
        return None


def ocr_image(img) -> str:
    if img.mode != "RGB":
        img = img.convert("RGB")
    return pytesseract.image_to_string(img, lang=OCR_LANG)


def extract_pdf(data: bytes) -> str:
    if not HAS_FITZ:
        return "(PyMuPDF nao instalado - pip install pymupdf)"
    doc = fitz.open(stream=data, filetype="pdf")
    out = []
    for i, page in enumerate(doc):
        t = page.get_text().strip()
        if len(t) < 10 and HAS_OCR:  # provavelmente pagina digitalizada
            pix = page.get_pixmap(dpi=200)
            t = ocr_image(Image.open(io.BytesIO(pix.tobytes("png")))).strip()
        if t:
            out.append(f"[pag {i + 1}]\n{t}")
    return "\n\n".join(out)


def extract_pptx(data: bytes) -> str:
    if not HAS_PPTX:
        return "(python-pptx nao instalado - pip install python-pptx)"
    prs = Presentation(io.BytesIO(data))
    out = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip():
                parts.append(shp.text_frame.text.strip())
            if getattr(shp, "has_table", False) and shp.has_table:
                for row in shp.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            if HAS_OCR and shp.shape_type == 13:  # PICTURE
                try:
                    o = ocr_image(Image.open(io.BytesIO(shp.image.blob))).strip()
                    if o:
                        parts.append("[imagem] " + o)
                except Exception:
                    pass
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append("[notas] " + slide.notes_slide.notes_text_frame.text.strip())
        if parts:
            out.append(f"[slide {i + 1}]\n" + "\n".join(parts))
    return "\n\n".join(out)


def extract_image(data: bytes, ext: str, query: str = "", use_vlm: bool = True) -> str:
    """OCR + interpretacao por VLM (quando disponivel) — combina os dois."""
    parts = []
    if use_vlm:
        mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{'png' if ext == 'png' else ext}"
        desc = describe_image_vlm(data, prompt=query, mime=mime)
        if desc:
            parts.append("[INTERPRETACAO VLM]\n" + desc)
    if HAS_OCR:
        try:
            txt = ocr_image(Image.open(io.BytesIO(data))).strip()
            if txt:
                parts.append("[TEXTO OCR]\n" + txt)
        except Exception:
            pass
    if not parts:
        return ("(Sem interpretacao: instala 'tesseract-ocr' para OCR e/ou define "
                "HF_TOKEN para o VLM gratuito.)")
    return "\n\n".join(parts)


IMG_EXTS = ("png", "jpg", "jpeg", "webp", "bmp", "tiff", "tif")


# ---------------------------------------------------------------------------
# Pesquisa web (DuckDuckGo, sem chave) + leitura de paginas
# ---------------------------------------------------------------------------
_TAG_RE = re.compile(r"<(script|style|noscript)[^>]*>.*?</\1>", re.S | re.I)
_HTML_RE = re.compile(r"<[^>]+>")


def _strip_html(html: str) -> str:
    html = _TAG_RE.sub(" ", html)
    text = _HTML_RE.sub(" ", html)
    text = re.sub(r"&nbsp;|&amp;|&lt;|&gt;|&quot;|&#\d+;", " ", text)
    return re.sub(r"\s{2,}", " ", text).strip()


def _fetch_page_text(url: str, cap: int = WEB_PAGE_CAP) -> str:
    try:
        r = requests.get(url, timeout=12, headers={
            "User-Agent": "Mozilla/5.0 (compatible; TICCopilot/1.0)"})
        r.raise_for_status()
        ctype = r.headers.get("content-type", "")
        if "html" not in ctype and "text" not in ctype:
            return ""
        return _strip_html(r.text)[:cap]
    except Exception:
        return ""


def web_search_context(query: str, max_results: int = WEB_MAX_RESULTS,
                       fetch_pages: int = WEB_FETCH_PAGES) -> tuple[str, list[dict]]:
    """Pesquisa DuckDuckGo + leitura das primeiras paginas. Devolve (bloco, fontes)."""
    if not HAS_DDG:
        return "", []
    try:
        with DDGS() as ddg:
            results = list(ddg.text(query, max_results=max_results, region="pt-pt"))
    except Exception:
        try:
            with DDGS() as ddg:
                results = list(ddg.text(query, max_results=max_results))
        except Exception:
            return "", []
    if not results:
        return "", []

    sources, parts = [], [
        "[RESULTADOS DA PESQUISA WEB - " + dt.datetime.now().strftime("%Y-%m-%d %H:%M") + "]",
        f"Pesquisa: {query}",
        "Usa estes resultados como evidencia atual; cita as fontes pelo numero [n].",
    ]
    used = 0
    for n, r in enumerate(results, 1):
        title = r.get("title", "")
        url = r.get("href") or r.get("url") or ""
        snippet = r.get("body", "")
        sources.append({"n": n, "title": title, "url": url})
        block = f"\n--- FONTE [{n}] {title}\nURL: {url}\nResumo: {snippet}"
        if n <= fetch_pages and url:
            page = _fetch_page_text(url)
            if page:
                block += f"\nConteudo da pagina:\n{page}"
        if used + len(block) > WEB_CONTEXT_CAP:
            break
        parts.append(block)
        used += len(block)
    parts.append("[FIM RESULTADOS WEB]")
    return "\n".join(parts), sources


# ---------------------------------------------------------------------------
# Speech-to-text local (faster-whisper, gratuito)
# ---------------------------------------------------------------------------
@st.cache_resource(show_spinner=False)
def load_whisper(size: str = WHISPER_DEFAULT):
    if not HAS_WHISPER:
        raise RuntimeError("faster-whisper nao instalado: pip install faster-whisper")
    return WhisperModel(size, device="cpu", compute_type="int8")


def transcribe_audio(audio_bytes: bytes, size: str = WHISPER_DEFAULT) -> str:
    model = load_whisper(size)
    segments, _info = model.transcribe(io.BytesIO(audio_bytes), language=None,
                                       vad_filter=True, beam_size=5)
    return " ".join(s.text.strip() for s in segments).strip()


# ---------------------------------------------------------------------------
# RAG HIBRIDO: BM25 + embeddings densos (FAISS) + RRF + reranker opcional
# ---------------------------------------------------------------------------
_WORD_RE = re.compile(r"[0-9a-zà-öø-ÿ_]+", re.I)


def _tokenize(text: str) -> list[str]:
    return _WORD_RE.findall((text or "").lower())


def _normalize_vecs(x):
    """Normaliza vetores para usar produto interno como similaridade cosseno."""
    if np is None:
        return x
    x = np.asarray(x, dtype="float32")
    if x.ndim == 1:
        x = x.reshape(1, -1)
    norms = np.linalg.norm(x, axis=1, keepdims=True)
    norms[norms == 0] = 1.0
    return x / norms


@st.cache_resource(show_spinner=False)
def load_embedding_model(model_name: str):
    """Carrega o modelo de embeddings uma vez por processo Streamlit."""
    if not HAS_SENTENCE_TRANSFORMERS:
        raise RuntimeError(
            "sentence-transformers nao esta instalado. Instala com: "
            "pip install sentence-transformers")
    return SentenceTransformer(model_name)


@st.cache_resource(show_spinner=False)
def load_reranker_model(model_name: str):
    """Carrega o reranker neural uma vez por processo Streamlit."""
    if not HAS_SENTENCE_TRANSFORMERS:
        raise RuntimeError(
            "sentence-transformers nao esta instalado. Instala com: "
            "pip install sentence-transformers")
    return CrossEncoder(model_name)


def _embed_documents(model, texts: list[str]):
    if hasattr(model, "encode_document"):
        try:
            return _normalize_vecs(model.encode_document(texts, convert_to_numpy=True, show_progress_bar=False))
        except TypeError:
            return _normalize_vecs(model.encode_document(texts))
    try:
        return _normalize_vecs(model.encode(texts, convert_to_numpy=True, normalize_embeddings=True, show_progress_bar=False))
    except TypeError:
        return _normalize_vecs(model.encode(texts, convert_to_numpy=True, show_progress_bar=False))


def _embed_query(model, query: str):
    if hasattr(model, "encode_query"):
        try:
            return _normalize_vecs(model.encode_query([query], convert_to_numpy=True, show_progress_bar=False))
        except TypeError:
            return _normalize_vecs(model.encode_query([query]))
    try:
        return _normalize_vecs(model.encode([query], convert_to_numpy=True, normalize_embeddings=True, show_progress_bar=False))
    except TypeError:
        return _normalize_vecs(model.encode([query], convert_to_numpy=True, show_progress_bar=False))


def _split_text_chunks(text: str, chunk_size: int = ATTACHMENT_CHUNK_SIZE,
                       overlap: int = ATTACHMENT_CHUNK_OVERLAP) -> list[tuple[int, int, str]]:
    """Divide texto em chunks com sobreposicao, devolvendo (inicio, fim, texto)."""
    text = (text or "").strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        if end < len(text):
            cut_candidates = [
                text.rfind("\n\n", start, end),
                text.rfind("\n", start, end),
                text.rfind(". ", start, end),
            ]
            cut = max(cut_candidates)
            if cut > start + chunk_size * 0.60:
                end = cut + 1
        chunk = text[start:end].strip()
        if chunk:
            chunks.append((start, end, chunk))
        if end >= len(text):
            break
        start = max(0, end - overlap)
        if start >= end:
            start = end
    return chunks


def _doc_hash(name: str, text: str) -> str:
    h = hashlib.sha256()
    h.update(name.encode("utf-8", "ignore"))
    h.update(b"\0")
    h.update((text or "").encode("utf-8", "ignore")[:200000])
    return h.hexdigest()[:16]


def _ensure_session_rag_docs():
    if "rag_docs" not in st.session_state:
        st.session_state.rag_docs = []
    if "rag_doc_hashes" not in st.session_state:
        st.session_state.rag_doc_hashes = set()


def add_docs_to_session_rag(docs: list[dict]) -> int:
    """Adiciona documentos extraidos ao indice logico da sessao, evitando duplicados."""
    _ensure_session_rag_docs()
    added = 0
    for doc in docs:
        dh = _doc_hash(doc["name"], doc["text"])
        if dh in st.session_state.rag_doc_hashes:
            continue
        d = dict(doc)
        d["hash"] = dh
        st.session_state.rag_docs.append(d)
        st.session_state.rag_doc_hashes.add(dh)
        added += 1
    return added


def clear_session_rag():
    st.session_state.rag_docs = []
    st.session_state.rag_doc_hashes = set()


def _build_faiss_index(embeddings):
    if not HAS_FAISS:
        raise RuntimeError("faiss-cpu nao esta instalado. Instala com: pip install faiss-cpu")
    dim = int(embeddings.shape[1])
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings.astype("float32"))
    return index


def _retrieve_attachment_chunks_hybrid(
    docs: list[dict],
    query: str,
    top_k: int,
    embedding_model_name: str,
    use_reranker: bool = False,
    reranker_model_name: str = RERANKER_MODEL_OPTIONS[0],
    candidate_multiplier: int = ATTACHMENT_CANDIDATE_MULTIPLIER,
) -> tuple[str, dict]:
    """RAG hibrido SOTA:

      chunks -> (1) BM25 lexical + (2) embeddings densos + FAISS
             -> fusao por Reciprocal Rank Fusion (RRF)
             -> reranker neural opcional (cross-encoder)
             -> top-k para o modelo
    Degrada graciosamente: so denso ou so BM25 se faltar uma das vias.
    """
    if not docs:
        return "", {"mode": "empty", "selected": [], "total_chunks": 0}

    candidates = []
    for doc_i, doc in enumerate(docs):
        for chunk_i, (start, end, chunk) in enumerate(_split_text_chunks(doc.get("text", ""))):
            if len(chunk) < 40:
                continue
            candidates.append({
                "doc_i": doc_i,
                "file": doc.get("name", f"doc_{doc_i}"),
                "hash": doc.get("hash", ""),
                "chunk_i": chunk_i,
                "start": start,
                "end": end,
                "text": chunk,
            })
    if not candidates:
        return "", {"mode": "empty", "selected": [], "total_chunks": 0}

    n_candidates = min(len(candidates), max(top_k, top_k * candidate_multiplier))
    dense_ranks: dict[int, int] = {}
    sparse_ranks: dict[int, int] = {}
    dense_ok = HAS_NUMPY and HAS_FAISS and HAS_SENTENCE_TRANSFORMERS
    sparse_ok = HAS_BM25

    # Via densa: embeddings + FAISS
    if dense_ok:
        try:
            model = load_embedding_model(embedding_model_name)
            doc_emb = _embed_documents(model, [c["text"] for c in candidates])
            query_emb = _embed_query(model, query)
            index = _build_faiss_index(doc_emb)
            scores, idxs = index.search(query_emb.astype("float32"), n_candidates)
            for rank, (idx, score) in enumerate(zip(idxs[0].tolist(), scores[0].tolist()), 1):
                if idx >= 0 and idx not in dense_ranks:
                    dense_ranks[idx] = rank
                    candidates[idx]["dense_score"] = float(score)
        except Exception as e:
            dense_ok = False
            for c in candidates:
                c["dense_error"] = str(e)[:120]

    # Via lexical: BM25 sobre os mesmos chunks
    if sparse_ok:
        try:
            tokenized = [_tokenize(c["text"]) for c in candidates]
            bm25 = BM25Okapi(tokenized)
            bm_scores = bm25.get_scores(_tokenize(query))
            order = sorted(range(len(candidates)), key=lambda i: bm_scores[i], reverse=True)
            for rank, idx in enumerate(order[:n_candidates], 1):
                if bm_scores[idx] <= 0:
                    break
                sparse_ranks[idx] = rank
                candidates[idx]["bm25_score"] = float(bm_scores[idx])
        except Exception:
            sparse_ok = False

    if not dense_ok and not sparse_ok:
        raise RuntimeError(
            "RAG indisponivel: instala 'sentence-transformers'+'faiss-cpu'+'numpy' "
            "(via densa) e/ou 'rank-bm25' (via lexical).")

    # Fusao por Reciprocal Rank Fusion
    fused: dict[int, float] = {}
    for idx, rank in dense_ranks.items():
        fused[idx] = fused.get(idx, 0.0) + 1.0 / (RRF_K + rank)
    for idx, rank in sparse_ranks.items():
        fused[idx] = fused.get(idx, 0.0) + 1.0 / (RRF_K + rank)

    ranked = []
    for idx, score in sorted(fused.items(), key=lambda kv: kv[1], reverse=True)[:n_candidates]:
        item = dict(candidates[idx])
        item["rrf_score"] = float(score)
        item["dense_rank"] = dense_ranks.get(idx)
        item["bm25_rank"] = sparse_ranks.get(idx)
        ranked.append(item)

    # Reranking neural opcional (cross-encoder): mais lento, mas mais preciso.
    if use_reranker and ranked and HAS_SENTENCE_TRANSFORMERS:
        try:
            reranker = load_reranker_model(reranker_model_name)
            pairs = [[query, r["text"]] for r in ranked]
            rr_scores = reranker.predict(pairs)
            for r, rr in zip(ranked, rr_scores):
                r["rerank_score"] = float(rr)
            ranked = sorted(ranked, key=lambda r: r.get("rerank_score", 0.0), reverse=True)
        except Exception as e:
            for r in ranked:
                r["rerank_error"] = str(e)[:120]

    selected = ranked[:top_k]

    via = []
    if dense_ok:
        via.append(f"densa ({embedding_model_name})")
    if sparse_ok:
        via.append("lexical (BM25)")
    parts = [
        "[RAG HIBRIDO EM ANEXOS]",
        "Os trechos abaixo foram recuperados dos anexos da sessao por pesquisa hibrida: "
        + " + ".join(via) + ", fundidos por Reciprocal Rank Fusion"
        + (f", reordenados por reranker neural ({reranker_model_name})" if use_reranker else "")
        + ".",
        "Usa estes trechos como evidencia. Se a resposta nao estiver nos trechos recuperados, diz isso explicitamente.",
        "Quando possivel, refere o ficheiro e o numero do trecho usado.",
        "[FIM NOTA RAG]",
    ]
    total = sum(len(d.get("text", "")) for d in docs)
    header = (f"Texto total indexado na sessao: {total:,} caracteres · "
              f"chunks avaliados: {len(candidates)} · fundidos por RRF: {len(ranked)} · "
              f"trechos enviados: {len(selected)}")
    parts.append(header.replace(",", " "))

    used_chars = 0
    for n, c in enumerate(selected, 1):
        score_bits = [f"rrf: {c.get('rrf_score', 0):.5f}"]
        if c.get("dense_rank"):
            score_bits.append("dense#%d (%.3f)" % (c["dense_rank"], c.get("dense_score", 0)))
        if c.get("bm25_rank"):
            score_bits.append("bm25#%d (%.2f)" % (c["bm25_rank"], c.get("bm25_score", 0)))
        if "rerank_score" in c:
            score_bits.append(f"rerank: {c['rerank_score']:.4f}")
        if "rerank_error" in c:
            score_bits.append(f"rerank_err: {c['rerank_error']}")
        block = (
            f"\n--- TRECHO RAG {n} | ficheiro: {c['file']} | "
            f"chunk: {c['chunk_i'] + 1} | chars: {c['start']}-{c['end']} | "
            + " | ".join(score_bits)
            + " ---\n"
            + c["text"]
        )
        if used_chars + len(block) > ATTACHMENT_CONTEXT_CAP:
            remaining = max(0, ATTACHMENT_CONTEXT_CAP - used_chars)
            if remaining > 1000:
                parts.append(block[:remaining] + "\n...[trecho cortado por limite de contexto]")
            break
        parts.append(block)
        used_chars += len(block)

    meta = {
        "mode": "hybrid_rag",
        "dense": dense_ok, "sparse": sparse_ok,
        "embedding_model": embedding_model_name if dense_ok else None,
        "reranker_model": reranker_model_name if use_reranker else None,
        "reranker_active": bool(use_reranker),
        "selected": [
            {k: c.get(k) for k in (
                "file", "chunk_i", "start", "end", "rrf_score", "dense_rank",
                "bm25_rank", "rerank_score")}
            for c in selected
        ],
        "total_chunks": len(candidates),
        "total_chars": total,
    }
    return "\n".join(parts).strip(), meta


def extract_files(
    files,
    query: str = "",
    use_attachment_search: bool = True,
    attachment_top_k: int = ATTACHMENT_TOP_K_DEFAULT,
    embedding_model_name: str = EMBEDDING_MODEL_OPTIONS[0],
    use_reranker: bool = False,
    reranker_model_name: str = RERANKER_MODEL_OPTIONS[0],
    use_vlm: bool = True,
) -> tuple[str, list[str], list[str], dict]:
    """Devolve (texto_contexto, nomes, avisos, meta).

    - Extrai texto de PDF/PPTX/TXT; imagens passam por VLM + OCR.
    - Guarda os documentos numa base RAG persistente na sessao (e na conversa).
    - Se o RAG estiver ativo, recupera os trechos mais relevantes (hibrido).
    - Se o RAG estiver inativo, usa o modo direto com limites/truncagem.
    """
    docs, names, warnings = [], [], []
    for f in files:
        name = f.name
        names.append(name)
        ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""
        data = f.getvalue()

        if ext == "pdf" and not HAS_FITZ:
            warnings.append(f"{name}: instala 'pymupdf' para ler PDF.")
            continue
        if ext in ("pptx", "ppt") and not HAS_PPTX:
            warnings.append(f"{name}: instala 'python-pptx' para ler PPTX.")
            continue

        try:
            if ext == "pdf":
                txt = extract_pdf(data)
            elif ext in ("pptx", "ppt"):
                txt = extract_pptx(data)
            elif ext in IMG_EXTS:
                txt = extract_image(data, ext, query=query, use_vlm=use_vlm)
            else:
                txt = data.decode("utf-8", "ignore")
        except Exception as e:
            warnings.append(f"{name}: falha a ler ({e}).")
            continue

        txt = (txt or "").strip()
        if not txt:
            warnings.append(f"{name}: sem texto extraido (PDF digitalizado? liga o OCR).")
            continue
        docs.append({"name": name, "text": txt})

    if not docs:
        return "", names, warnings, {"mode": "empty"}

    added = add_docs_to_session_rag(docs)
    if added:
        warnings.append(f"{added} documento(s) adicionados ao indice RAG da conversa.")

    raw_total = sum(len(d["text"]) for d in docs)

    if use_attachment_search:
        ctx, meta = _retrieve_attachment_chunks_hybrid(
            st.session_state.rag_docs,
            query,
            max(1, attachment_top_k),
            embedding_model_name,
            use_reranker=use_reranker,
            reranker_model_name=reranker_model_name,
        )
        return ctx, names, warnings, meta

    # Modo direto/truncado.
    chunks = []
    for doc in docs:
        txt = doc["text"]
        if len(txt) > PER_FILE_CAP:
            txt = txt[:PER_FILE_CAP] + "\n...[ficheiro truncado]"
        chunks.append(f"--- {doc['name']} ---\n{txt}")

    full = "\n\n".join(chunks)
    if len(full) > TOTAL_CAP:
        full = full[:TOTAL_CAP] + "\n...[total truncado]"
    return full, names, warnings, {"mode": "direct", "total_chars": raw_total}


def retrieve_session_rag_context(
    query: str,
    attachment_top_k: int = ATTACHMENT_TOP_K_DEFAULT,
    embedding_model_name: str = EMBEDDING_MODEL_OPTIONS[0],
    use_reranker: bool = False,
    reranker_model_name: str = RERANKER_MODEL_OPTIONS[0],
) -> tuple[str, dict]:
    """Recupera contexto da base RAG da conversa mesmo sem novos anexos."""
    _ensure_session_rag_docs()
    if not st.session_state.rag_docs:
        return "", {"mode": "empty"}
    return _retrieve_attachment_chunks_hybrid(
        st.session_state.rag_docs,
        query,
        max(1, attachment_top_k),
        embedding_model_name,
        use_reranker=use_reranker,
        reranker_model_name=reranker_model_name,
    )


# ---------------------------------------------------------------------------
# Chaves de API
# ---------------------------------------------------------------------------
def api_key_for(model_key: str) -> str:
    cfg = MODELS[model_key]
    env_value = os.getenv(cfg["env"])
    if env_value:
        return env_value.strip()
    try:
        if cfg["env"] in st.secrets:
            return str(st.secrets[cfg["env"]]).strip()
    except Exception:
        pass
    st.error(f"Secret em falta: configura `{cfg['env']}` no ambiente de deploy.")
    st.stop()


def new_thread_id() -> str:
    return secrets.token_urlsafe(16)


# ---------------------------------------------------------------------------
# Rate limiting: retry com backoff + throttle global
# ---------------------------------------------------------------------------
RATE_LIMIT_RETRIES = 4
RATE_LIMIT_BASE_DELAY = 2.0
# Chunk pequeno para o requests entregar eventos assim que chegam.
STREAM_CHUNK_SIZE = 64
_min_interval = [1.2]
_throttle_lock = threading.Lock()
_last_request_ts = [0.0]


def _throttle():
    with _throttle_lock:
        wait = _min_interval[0] - (time.time() - _last_request_ts[0])
        if wait > 0:
            time.sleep(wait)
        _last_request_ts[0] = time.time()


def _is_rate_limit(text) -> bool:
    t = str(text).lower()
    return any(s in t for s in ("429", "rate limit", "rate-limit", "too many requests"))


# ---------------------------------------------------------------------------
# Streaming do agente IAedu
# ---------------------------------------------------------------------------
def _decode_stream(cfg: dict, api_key: str, thread_id: str, message: str):
    fields = {
        "channel_id": (None, cfg["channel_id"]),
        "thread_id": (None, thread_id),
        "user_info": (None, "{}"),
        "message": (None, message),
    }
    headers = {"x-api-key": api_key, "Accept": "text/event-stream"}
    decoder = json.JSONDecoder()
    buffer = ""
    got_tokens = False
    fallback_msg = None

    _throttle()
    with requests.post(cfg["endpoint"], headers=headers, files=fields,
                       stream=True, timeout=300) as resp:
        resp.raise_for_status()
        for raw in resp.iter_content(chunk_size=STREAM_CHUNK_SIZE, decode_unicode=True):
            if not raw:
                continue
            buffer += raw
            while True:
                buffer = buffer.lstrip()
                if not buffer:
                    break
                try:
                    event, idx = decoder.raw_decode(buffer)
                except json.JSONDecodeError:
                    break
                buffer = buffer[idx:]
                if not isinstance(event, dict):
                    continue
                etype = event.get("type")
                if etype == "token":
                    tok = event.get("content", "")
                    if tok:
                        got_tokens = True
                        yield ("token", tok)
                elif etype == "message":
                    content = event.get("content", {})
                    if isinstance(content, dict):
                        fallback_msg = content.get("content")
                elif etype == "context_limit":
                    yield ("context_limit", None)
                    return
                elif etype == "error":
                    yield ("error", str(event.get("content", "erro desconhecido")))
                    return

    if not got_tokens and fallback_msg:
        yield ("token", fallback_msg)


def respond(model_key: str, message: str, thread_id: str, api_key: str, persist: bool = False):
    cfg = MODELS[model_key]
    tid = thread_id
    ctx_retried = False
    rl_attempts = 0
    while True:
        yielded_any = False
        hit_ctx = False
        rate_limited = False
        try:
            for kind, payload in _decode_stream(cfg, api_key, tid, message):
                if kind == "token":
                    yielded_any = True
                    yield payload
                elif kind == "context_limit":
                    hit_ctx = True
                    tid = new_thread_id()
                    if persist:
                        st.session_state.threads[model_key] = tid
                    break
                elif kind == "error":
                    if _is_rate_limit(payload) and not yielded_any:
                        rate_limited = True
                        break
                    yield f"\n\n> Erro do agente: {payload}"
                    return
        except requests.HTTPError as e:
            status = getattr(e.response, "status_code", None)
            if status == 429 and not yielded_any:
                rate_limited = True
            else:
                yield f"\n\n> Erro HTTP ({cfg['label']}): {e}"
                return
        except Exception as e:
            yield f"\n\n> Falha de ligacao ({cfg['label']}): {e}"
            return

        if rate_limited:
            if rl_attempts < RATE_LIMIT_RETRIES:
                rl_attempts += 1
                delay = RATE_LIMIT_BASE_DELAY * (2 ** (rl_attempts - 1))
                yield (f"\n\n_(limite de pedidos atingido - nova tentativa em "
                       f"{delay:.0f}s . {rl_attempts}/{RATE_LIMIT_RETRIES})_\n\n")
                time.sleep(delay)
                continue
            yield (f"\n\n> Limite de pedidos do IAedu (429) persistente apos "
                   f"{RATE_LIMIT_RETRIES} tentativas. Espera um pouco e tenta de novo.")
            return

        if hit_ctx and not ctx_retried:
            ctx_retried = True
            yield "\n\n_(contexto cheio - recomecei a conversa deste modelo)_\n\n"
            continue
        return


# ---------------------------------------------------------------------------
# Estimativas e contabilizacao de custos da sessao
# ---------------------------------------------------------------------------
def est_tokens(text: str) -> int:
    return max(1, len(text) // 4)


def est_cost(prompt: str, output: str, model_key: str) -> float:
    cfg = MODELS[model_key]
    return (est_tokens(prompt) / 1_000_000 * cfg["price_in"]
            + est_tokens(output) / 1_000_000 * cfg["price_out"])


def track_usage(model_key: str, prompt: str, output: str):
    u = st.session_state.setdefault("usage", {"cost": 0.0, "calls": 0, "tokens_out": 0})
    u["cost"] += est_cost(prompt, output, model_key)
    u["calls"] += 1
    u["tokens_out"] += est_tokens(output)


# ---------------------------------------------------------------------------
# Streaming suave: efeito maquina-de-escrever consistente
# ---------------------------------------------------------------------------
def smooth_stream(gen, chunk_chars: int = 3, delay: float = 0.010):
    """Reparte blocos grandes em micro-pedacos com um pequeno atraso.

    Alguns agentes/proxies entregam o texto em blocos grandes (ou tudo no fim);
    isto garante que a resposta aparece SEMPRE em streaming fluido, caracter a
    caracter, sem alterar o conteudo.
    """
    for piece in gen:
        if not piece:
            continue
        if len(piece) <= chunk_chars * 2:
            yield piece
            continue
        for i in range(0, len(piece), chunk_chars):
            yield piece[i:i + chunk_chars]
            if delay:
                time.sleep(delay)


# ---------------------------------------------------------------------------
# Composicao da mensagem (memoria + conversa + anexos + web + pergunta)
# ---------------------------------------------------------------------------
def _compact_for_context(text: str, cap: int = 6000) -> str:
    """Limita blocos muito longos sem rebentar o contexto."""
    text = (text or "").strip()
    if len(text) <= cap:
        return text
    head = cap // 2
    tail = cap - head
    return text[:head] + "\n...[texto intermedio omitido]...\n" + text[-tail:]


def _assistant_msg_to_text(msg: dict) -> str:
    """Converte os varios modos da UI para texto utilizavel como contexto."""
    mode = msg.get("mode")
    if mode == "single":
        model = MODELS.get(msg.get("model", ""), {}).get("label", msg.get("model", "assistant"))
        return f"Assistente ({model}):\n{msg.get('content', '')}"

    if mode == "duel":
        data = msg.get("data", {})
        return (
            "Assistente (Duelo):\n"
            f"[Claude]\n{data.get('claude', '')}\n\n"
            f"[GPT]\n{data.get('gpt', '')}"
        )

    if mode == "collab":
        ex = MODELS.get(msg.get("executor", ""), {}).get("label", msg.get("executor", "executor"))
        rv = MODELS.get(msg.get("reviewer", ""), {}).get("label", msg.get("reviewer", "reviewer"))
        parts = [
            f"Assistente (Colaboracao; executor={ex}; revisor={rv}):",
            f"[Rascunho]\n{msg.get('draft', '')}",
            f"[Revisao]\n{msg.get('critique', '')}",
        ]
        if msg.get("synthesis"):
            parts.append(f"[Versao final]\n{msg.get('synthesis', '')}")
        return "\n\n".join(parts)

    return f"Assistente:\n{msg.get('content', '')}"


def session_transcript(max_chars: int = SESSION_TRANSCRIPT_CAP) -> str:
    """Transcricao da conversa atual, em janela deslizante (fim -> inicio)."""
    messages = st.session_state.get("messages", [])
    blocks = []

    for msg in messages:
        if msg.get("role") == "user":
            extra = ""
            if msg.get("attachments"):
                extra = "\n[anexos neste turno: " + ", ".join(msg["attachments"]) + "]"
            blocks.append("Utilizador:\n" + msg.get("content", "") + extra)
        elif msg.get("role") == "assistant":
            blocks.append(_assistant_msg_to_text(msg))

    selected = []
    total = 0
    for block in reversed(blocks):
        block = _compact_for_context(block, cap=6000)
        add = len(block) + 2
        if selected and total + add > max_chars:
            break
        selected.append(block)
        total += add

    selected.reverse()
    return "\n\n---\n\n".join(selected).strip()


def build_agent_message(
    user_prompt: str,
    username: str,
    use_memory: bool,
    attach_text: str,
    use_session_context: bool = True,
    web_text: str = "",
) -> str:
    parts = []

    if use_memory:
        mem = load_memory(username).get("master", "").strip()
        if mem:
            parts.append("[MEMORIA DO UTILIZADOR - contexto de fundo, nao e a pergunta]\n"
                         + mem + "\n[FIM MEMORIA]")

    if use_session_context:
        hist = session_transcript()
        if hist:
            parts.append(
                "[TRANSCRICAO DA CONVERSA ATUAL - usa para manter continuidade; "
                "nao repitas a menos que seja util]\n"
                + hist +
                "\n[FIM TRANSCRICAO]"
            )

    if attach_text:
        parts.append("[FICHEIROS / ANEXOS DA CONVERSA]\n" + attach_text + "\n[FIM FICHEIROS]")

    if web_text:
        parts.append(web_text)

    parts.append("[PERGUNTA ATUAL]\n" + user_prompt)
    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Saudacao por hora do dia (hero personalizado)
# ---------------------------------------------------------------------------
def greeting_pt(name: str = "") -> str:
    h = dt.datetime.now().hour
    if 6 <= h < 13:
        g = "Bom dia"
    elif 13 <= h < 20:
        g = "Boa tarde"
    else:
        g = "Boa noite"
    first = (name or "").strip().split()[0] if (name or "").strip() else ""
    return f"{g}, {first}" if first else g


def date_pt() -> str:
    wd = ["segunda-feira", "ter\u00e7a-feira", "quarta-feira", "quinta-feira",
          "sexta-feira", "s\u00e1bado", "domingo"]
    ms = ["janeiro", "fevereiro", "mar\u00e7o", "abril", "maio", "junho", "julho",
          "agosto", "setembro", "outubro", "novembro", "dezembro"]
    n = dt.datetime.now()
    return f"{wd[n.weekday()]}, {n.day} de {ms[n.month - 1]}"


# ---------------------------------------------------------------------------
# Tema / CSS — "blueprint industrial" (TIC: CAD, 3D printing, 3D reality)
# ---------------------------------------------------------------------------
THEME_CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;500;600;700&family=Inter:wght@400;500;600&family=JetBrains+Mono:wght@400;500;600&display=swap');
:root{
  --bg:#090D13; --bg-2:#0F141D; --bg-3:#141B27;
  --line:#1F2937; --line-2:#2C3A52;
  --text:#E8EDF5; --muted:#8B96AB; --faint:#5A6679;
  --claude:#D2A24C; --gpt:#46B98C; --tic:#5CC8E8; --accent:#D2A24C;
  --font-display:'Space Grotesk',system-ui,sans-serif;
  --font-ui:'Inter',system-ui,sans-serif;
  --font-mono:'JetBrains Mono',ui-monospace,monospace;
}
.stApp,[data-testid="stAppViewContainer"]{
  background:
    /* grelha blueprint muito subtil */
    repeating-linear-gradient(0deg, rgba(92,200,232,0.022) 0 1px, transparent 1px 56px),
    repeating-linear-gradient(90deg, rgba(92,200,232,0.022) 0 1px, transparent 1px 56px),
    radial-gradient(1100px 560px at 85% -12%, rgba(210,162,76,0.06), transparent 60%),
    radial-gradient(900px 520px at -8% 4%, rgba(92,200,232,0.05), transparent 55%),
    var(--bg);
  color:var(--text); font-family:var(--font-ui);
}
[data-testid="stHeader"]{ background:transparent; }
[data-testid="stMainBlockContainer"],.block-container{ max-width:1080px; padding-top:1.1rem; padding-bottom:5rem; }
body,p,li,span,div{ font-family:var(--font-ui); }
h1,h2,h3{ font-family:var(--font-display); letter-spacing:-.01em; }
[data-testid="stSidebar"]{ background:linear-gradient(180deg,var(--bg-2),#0C111A); border-right:1px solid var(--line); }
[data-testid="stSidebar"] .block-container{ padding-top:1.3rem; }

/* eyebrows com marca de cota (dimension line) */
.eyebrow{ font-family:var(--font-mono); text-transform:uppercase; letter-spacing:.22em; font-size:.62rem;
  color:var(--faint); margin:1.15rem 0 .55rem; display:flex; align-items:center; gap:.6rem; }
.eyebrow::before{ content:""; width:7px; height:7px; border:1px solid var(--line-2); transform:rotate(45deg); flex:none; }
.eyebrow::after{ content:""; height:1px; flex:1;
  background:linear-gradient(90deg,var(--line-2),transparent); }

[data-testid="stChatMessage"]{ background:var(--bg-2); border:1px solid var(--line); border-radius:16px;
  padding:1.05rem 1.25rem; box-shadow:0 1px 0 rgba(255,255,255,.02),0 14px 34px rgba(0,0,0,.32);
  position:relative; }
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]){
  background:linear-gradient(180deg, rgba(92,200,232,.035), transparent 70%);
  border:1px dashed var(--line-2); box-shadow:none; }
[data-testid^="stChatMessageAvatar"]{ display:none; }
[data-testid="stChatMessage"] > div:first-child{ gap:0 !important; }
[data-testid="stChatMessage"] p{ line-height:1.64; }
[data-testid="stChatMessage"] code{ font-family:var(--font-mono); background:#0b1018; border:1px solid var(--line);
  padding:.05rem .35rem; border-radius:6px; font-size:.86em; }
[data-testid="stChatMessage"] pre{ background:#0b1018 !important; border:1px solid var(--line); border-radius:12px; }
a{ color:var(--tic); }

[data-testid="stChatInput"]{ background:var(--bg-2); border:1px solid var(--line); border-radius:14px; }
[data-testid="stChatInput"]:focus-within{ border-color:var(--tic); box-shadow:0 0 0 3px rgba(92,200,232,.13); }

.stButton > button{ background:transparent; color:var(--text); border:1px solid var(--line-2); border-radius:10px;
  font-family:var(--font-mono); text-transform:uppercase; letter-spacing:.09em; font-size:.7rem;
  padding:.5rem .85rem; transition:all .15s ease; }
.stButton > button:hover{ border-color:var(--tic); color:var(--tic); background:rgba(92,200,232,.06); }
.stDownloadButton > button{ background:transparent; color:var(--text); border:1px solid var(--line-2);
  border-radius:10px; font-family:var(--font-mono); font-size:.7rem; letter-spacing:.09em;
  text-transform:uppercase; }
.stDownloadButton > button:hover{ border-color:var(--accent); color:var(--accent); }

[data-testid="stExpander"]{ border:1px solid var(--line); border-radius:12px; background:var(--bg-2); }
[data-testid="stExpander"] summary{ font-family:var(--font-mono); font-size:.78rem; letter-spacing:.04em; }
[data-testid="stVerticalBlockBorderWrapper"]{ background:var(--bg-3); border-radius:14px; }
[data-testid="stFileUploaderDropzone"]{ background:var(--bg-3); border:1px dashed var(--line-2); border-radius:12px; }
[data-testid="stPopover"] button{ font-family:var(--font-mono); }

.badge{ display:inline-flex; align-items:center; gap:.55rem; font-family:var(--font-mono); font-size:.72rem;
  padding:.32rem .7rem; border:1px solid var(--line-2); border-radius:999px; background:var(--bg-3); color:var(--text); }
.badge .dot{ width:8px; height:8px; border-radius:50%; box-shadow:0 0 10px 0 currentColor; }
.badge--claude{ border-color:rgba(210,162,76,.4); } .badge--claude .dot{ background:var(--claude); color:var(--claude); }
.badge--gpt{ border-color:rgba(70,185,140,.4); } .badge--gpt .dot{ background:var(--gpt); color:var(--gpt); }
.badge--tic{ border-color:rgba(92,200,232,.4); } .badge--tic .dot{ background:var(--tic); color:var(--tic); }
.badge .sub{ color:var(--faint); margin-left:.35rem; }

.reason{ margin:.55rem 0 .2rem; padding:.6rem .85rem; border-left:2px solid var(--line-2); background:var(--bg-3);
  border-radius:0 10px 10px 0; color:var(--muted); font-size:.85rem; line-height:1.5; }
.reason b{ color:var(--text); font-weight:600; }
.reason--claude{ border-left-color:var(--claude); } .reason--gpt{ border-left-color:var(--gpt); }

.metricline{ font-family:var(--font-mono); font-size:.7rem; color:var(--muted); display:flex; gap:.55rem;
  flex-wrap:wrap; margin-top:.5rem; } .metricline .sep{ color:var(--faint); }

.step{ font-family:var(--font-mono); text-transform:uppercase; letter-spacing:.14em; font-size:.7rem;
  color:var(--muted); margin:1.1rem 0 .35rem; display:flex; align-items:center; gap:.6rem; }
.step .n{ width:1.35rem; height:1.35rem; border-radius:6px; background:var(--text); color:var(--bg);
  display:inline-flex; align-items:center; justify-content:center; font-weight:600; font-size:.72rem; }

.col-head{ display:flex; align-items:center; justify-content:space-between; margin:.2rem 0 .5rem; }
.routerule{ display:flex; gap:.55rem; align-items:flex-start; padding:.35rem 0; font-size:.8rem;
  color:var(--text); border-bottom:1px solid var(--line); }
.routerule:last-child{ border-bottom:none; }
.routerule .d{ width:8px; height:8px; border-radius:50%; margin-top:.42rem; flex:none; }
.routerule .b{ display:block; font-family:var(--font-mono); font-size:.66rem; color:var(--faint); margin-top:.1rem; }
.attach-note{ font-family:var(--font-mono); font-size:.68rem; color:var(--faint); margin-top:.4rem; }
.src-note{ font-family:var(--font-mono); font-size:.68rem; color:var(--faint); margin-top:.45rem; line-height:1.7; }
.src-note a{ color:var(--tic); text-decoration:none; border-bottom:1px dotted var(--line-2); }

.userbox{ display:flex; align-items:center; gap:.6rem; padding:.55rem .7rem; border:1px solid var(--line);
  border-radius:12px; background:var(--bg-3); margin-bottom:.4rem; }
.userbox .av{ width:30px; height:30px; border-radius:8px; background:linear-gradient(135deg,var(--tic),var(--claude));
  display:flex; align-items:center; justify-content:center; font-family:var(--font-mono); font-weight:600; color:#0A0E14; }
.userbox .n{ font-size:.85rem; } .userbox .r{ font-family:var(--font-mono); font-size:.62rem; color:var(--faint); }

.convtitle{ font-family:var(--font-mono); font-size:.68rem; color:var(--faint); letter-spacing:.06em; }
.usagebox{ font-family:var(--font-mono); font-size:.68rem; color:var(--muted); border:1px solid var(--line);
  border-radius:10px; padding:.5rem .7rem; background:var(--bg-3); display:flex; gap:.8rem; }
.usagebox b{ color:var(--tic); font-weight:600; }

.authwrap{ max-width:420px; margin:.5rem auto 0; }
.authwrap h3{ font-family:var(--font-display); font-weight:600; font-size:1.4rem; margin-bottom:.2rem; }
.authwrap p.s{ color:var(--muted); font-size:.85rem; margin-bottom:.4rem; }

/* ------- efeitos ------- */
@keyframes msgIn{ from{opacity:0; transform:translateY(10px) scale(.992);} to{opacity:1; transform:none;} }
[data-testid="stChatMessage"]{ animation:msgIn .42s cubic-bezier(.2,.7,.25,1); transition:border-color .2s ease, transform .2s ease; }
[data-testid="stChatMessage"]:hover{ border-color:var(--line-2); }
@keyframes dotPulse{ 0%,100%{ box-shadow:0 0 6px 0 currentColor; opacity:.95; } 50%{ box-shadow:0 0 14px 2px currentColor; opacity:1; } }
.badge .dot{ animation:dotPulse 2.4s ease-in-out infinite; }
@keyframes reasonSweep{ from{ background-position:-200% 0; } to{ background-position:200% 0; } }
.reason{ background-image:linear-gradient(100deg, transparent 35%, rgba(255,255,255,.025) 50%, transparent 65%), linear-gradient(var(--bg-3),var(--bg-3)); background-size:200% 100%, 100% 100%; animation:reasonSweep 5s linear infinite; }
[data-testid="stStatusWidget"], .stSpinner > div{ color:var(--tic) !important; }
@keyframes scanIn{ from{ transform:scaleX(0); } to{ transform:scaleX(1); } }
.step::after{ content:""; height:1px; flex:1; background:linear-gradient(90deg,var(--line-2),transparent); transform-origin:left; animation:scanIn .6s ease; }
::selection{ background:rgba(92,200,232,.28); }
*::-webkit-scrollbar{ width:10px; height:10px; }
*::-webkit-scrollbar-thumb{ background:var(--line-2); border-radius:8px; border:2px solid var(--bg); }
*::-webkit-scrollbar-track{ background:transparent; }
</style>
"""

# ---------------------------------------------------------------------------
# Hero — peca 3D wireframe a rodar sobre grelha blueprint
# ---------------------------------------------------------------------------
HERO_HTML = """
<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
@import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@500;600;700&family=Inter:wght@400;500&family=JetBrains+Mono:wght@400;500&display=swap');
*{margin:0;padding:0;box-sizing:border-box;}
html,body{height:100%;overflow:hidden;background:#090D13;}
#c{position:absolute;inset:0;}
.wrap{position:relative;height:100%;display:flex;flex-direction:column;justify-content:center;
  padding:26px 32px;font-family:'Inter',sans-serif;color:#E8EDF5;}
.eye{font-family:'JetBrains Mono',monospace;text-transform:uppercase;letter-spacing:.32em;
  font-size:10.5px;color:#5A6679;margin-bottom:12px;display:flex;align-items:center;gap:10px;}
.eye .sq{width:7px;height:7px;border:1px solid #2C3A52;transform:rotate(45deg);}
.greet{font-family:'JetBrains Mono',monospace;font-size:12.5px;color:#5CC8E8;letter-spacing:.1em;
  margin-bottom:9px;opacity:0;animation:greetIn .9s ease .15s forwards;}
.greet .dt{color:#5A6679;}
@keyframes greetIn{from{opacity:0;transform:translateY(6px);}to{opacity:1;transform:none;}}
h1{font-family:'Space Grotesk',sans-serif;font-weight:700;font-size:46px;line-height:1;letter-spacing:-1px;}
h1 .tic{color:#5CC8E8;}
h1 .x{color:#D2A24C;margin:0 8px;font-weight:500;}
.sub{margin-top:14px;font-family:'JetBrains Mono',monospace;font-size:13px;color:#8B96AB;height:18px;}
.sub .cur{color:#5CC8E8;}
.pills{position:absolute;top:24px;right:30px;display:flex;gap:8px;}
.pill{font-family:'JetBrains Mono',monospace;font-size:11px;color:#E8EDF5;border:1px solid #2C3A52;
  border-radius:999px;padding:5px 11px;display:flex;align-items:center;gap:7px;background:rgba(20,27,39,.65);}
.pill .d{width:7px;height:7px;border-radius:50%;}
.corner{position:absolute;width:16px;height:16px;border:1px solid #2C3A52;opacity:.8;}
.corner.tl{top:12px;left:12px;border-right:none;border-bottom:none;}
.corner.br{bottom:12px;right:12px;border-left:none;border-top:none;}
.line{position:absolute;left:0;right:0;bottom:0;height:1px;
  background:linear-gradient(90deg,transparent,#5CC8E855,#D2A24C55,#46B98C55,transparent);}
</style></head><body>
<canvas id="c"></canvas>
<div class="corner tl"></div><div class="corner br"></div>
<div class="pills">
  <div class="pill"><span class="d" style="background:#D2A24C;box-shadow:0 0 8px #D2A24C"></span>Opus 4.7</div>
  <div class="pill"><span class="d" style="background:#46B98C;box-shadow:0 0 8px #46B98C"></span>GPT-5.5</div>
</div>
<div class="wrap">
  <div class="eye"><span class="sq"></span>Autoeuropa &middot; Technical Innovation Center</div>
  <div class="greet">__GREETING__</div>
  <h1><span class="tic">TIC</span> Copilot <span class="x">&middot;</span> Claude <span class="x">&times;</span> GPT</h1>
  <div class="sub"><span id="t"></span><span class="cur">_</span></div>
</div>
<div class="line"></div>
<script>
var canvas=document.getElementById('c'),ctx=canvas.getContext('2d'),W,H;
function size(){W=canvas.width=canvas.offsetWidth;H=canvas.height=canvas.offsetHeight;}
size();window.addEventListener('resize',size);

// Perfil de engrenagem (12 dentes) extrudido — wireframe CAD
var TEETH=12, prof=[];
for(var i=0;i<TEETH*4;i++){
  var a=i/(TEETH*4)*Math.PI*2;
  var step=i%4, r=(step===0||step===3)?0.72:1.0;
  prof.push([Math.cos(a)*r, Math.sin(a)*r]);
}
var bore=[]; for(var i=0;i<24;i++){var a=i/24*Math.PI*2; bore.push([Math.cos(a)*0.32, Math.sin(a)*0.32]);}
function proj(p, rx, ry, s, cx, cy){
  var x=p[0], y=p[1], z=p[2];
  var c1=Math.cos(ry), s1=Math.sin(ry);     // rot Y
  var x1=x*c1+z*s1, z1=-x*s1+z*c1;
  var c2=Math.cos(rx), s2=Math.sin(rx);     // rot X
  var y2=y*c2-z1*s2, z2=y*s2+z1*c2;
  var d=3.2/(3.2+z2);
  return [cx+x1*s*d, cy+y2*s*d];
}
function ring(pts, z){ return pts.map(function(p){return [p[0],p[1],z];}); }
function drawLoop(pts3, rx, ry, s, cx, cy, color){
  ctx.strokeStyle=color; ctx.lineWidth=1; ctx.beginPath();
  for(var i=0;i<=pts3.length;i++){
    var q=proj(pts3[i%pts3.length], rx, ry, s, cx, cy);
    if(i===0)ctx.moveTo(q[0],q[1]); else ctx.lineTo(q[0],q[1]);
  }
  ctx.stroke();
}
var t0=Date.now();
function draw(){
  ctx.clearRect(0,0,W,H);
  var t=(Date.now()-t0)/1000;
  var rx=0.42+Math.sin(t*0.21)*0.10, ry=t*0.35;
  var cx=W*0.80, cy=H*0.52, s=Math.min(H*0.62, 120), dz=0.22;
  var front=ring(prof,dz), back=ring(prof,-dz);
  // arestas laterais
  ctx.strokeStyle='rgba(92,200,232,0.20)'; ctx.lineWidth=1; ctx.beginPath();
  for(var i=0;i<prof.length;i+=2){
    var a=proj(front[i],rx,ry,s,cx,cy), b=proj(back[i],rx,ry,s,cx,cy);
    ctx.moveTo(a[0],a[1]); ctx.lineTo(b[0],b[1]);
  }
  ctx.stroke();
  drawLoop(front,rx,ry,s,cx,cy,'rgba(92,200,232,0.55)');
  drawLoop(back,rx,ry,s,cx,cy,'rgba(92,200,232,0.28)');
  drawLoop(ring(bore,dz),rx,ry,s,cx,cy,'rgba(210,162,76,0.55)');
  drawLoop(ring(bore,-dz),rx,ry,s,cx,cy,'rgba(210,162,76,0.30)');
  // nos dos dentes (frente)
  for(var i=0;i<prof.length;i+=4){
    var q=proj(front[i],rx,ry,s,cx,cy);
    ctx.fillStyle='rgba(70,185,140,0.75)';
    ctx.beginPath(); ctx.arc(q[0],q[1],1.6,0,6.3); ctx.fill();
  }
  // cruz de centro (estilo desenho tecnico)
  var c=proj([0,0,0],rx,ry,s,cx,cy);
  ctx.strokeStyle='rgba(139,150,171,0.35)';
  ctx.beginPath();
  ctx.moveTo(c[0]-10,c[1]); ctx.lineTo(c[0]+10,c[1]);
  ctx.moveTo(c[0],c[1]-10); ctx.lineTo(c[0],c[1]+10);
  ctx.stroke();
  requestAnimationFrame(draw);
}
draw();
var phrases=["Gestao de pedidos CAD & impressao 3D","Roteamento inteligente por tarefa",
  "RAG hibrido sobre os teus documentos","Pesquisa web, voz e visao por computador",
  "Historico de conversas sempre disponivel"];
var pi=0,ci=0,del=false,el=document.getElementById('t');
function type(){var w=phrases[pi];el.textContent=w.substring(0,ci);
  if(!del){ci++;if(ci>w.length){del=true;setTimeout(type,1600);return;}}
  else{ci--;if(ci<0){del=false;pi=(pi+1)%phrases.length;ci=0;}}
  setTimeout(type,del?26:52);}
type();
</script></body></html>
"""


def hero_html(greeting: str, date_str: str = "") -> str:
    """Hero personalizado: 'Bom dia, Carlos' consoante a hora do dia."""
    line = greeting + (f" <span class='dt'>&middot; {date_str}</span>" if date_str else "")
    return HERO_HTML.replace("__GREETING__", line)


# ---------------------------------------------------------------------------
# Pagina de login: fundo aurora animado + particulas que formam "TIC COPILOT"
# ---------------------------------------------------------------------------
LOGIN_CSS = """
<style>
@keyframes aurora{ 0%{background-position:0% 50%;} 50%{background-position:100% 50%;} 100%{background-position:0% 50%;} }
.stApp,[data-testid="stAppViewContainer"]{
  background:linear-gradient(115deg,#090D13 0%,#0A1620 22%,#10222E 42%,#1B1610 62%,#0C1A22 82%,#090D13 100%) !important;
  background-size:320% 320% !important;
  animation:aurora 22s ease-in-out infinite;
}
@keyframes gridDrift{ from{background-position:0 0,0 0;} to{background-position:56px 56px,56px 56px;} }
.stApp::before{
  content:""; position:fixed; inset:0; pointer-events:none; z-index:0;
  background:
    repeating-linear-gradient(0deg, rgba(92,200,232,0.030) 0 1px, transparent 1px 56px),
    repeating-linear-gradient(90deg, rgba(92,200,232,0.030) 0 1px, transparent 1px 56px);
  animation:gridDrift 26s linear infinite;
  mask-image:radial-gradient(900px 600px at 50% 36%, #000 30%, transparent 80%);
  -webkit-mask-image:radial-gradient(900px 600px at 50% 36%, #000 30%, transparent 80%);
}
/* cartao de login em vidro */
[data-testid="stForm"]{
  background:rgba(13,18,27,.55) !important;
  backdrop-filter:blur(16px) saturate(130%);
  -webkit-backdrop-filter:blur(16px) saturate(130%);
  border:1px solid rgba(92,200,232,.22) !important;
  border-radius:18px !important;
  padding:1.35rem 1.35rem 1.1rem !important;
  box-shadow:0 30px 70px rgba(0,0,0,.55), inset 0 1px 0 rgba(255,255,255,.05);
  transition:border-color .25s ease, box-shadow .25s ease;
}
[data-testid="stForm"]:focus-within{
  border-color:rgba(92,200,232,.55) !important;
  box-shadow:0 30px 70px rgba(0,0,0,.55), 0 0 0 4px rgba(92,200,232,.10), inset 0 1px 0 rgba(255,255,255,.06);
}
[data-testid="stForm"] .stButton > button,[data-testid="stForm"] button[kind="primaryFormSubmit"],
[data-testid="stForm"] button[kind="secondaryFormSubmit"]{
  background:linear-gradient(120deg, rgba(92,200,232,.16), rgba(210,162,76,.14)) !important;
  border:1px solid rgba(92,200,232,.4) !important; color:#E8EDF5 !important;
}
[data-testid="stForm"] button:hover{ border-color:#5CC8E8 !important; box-shadow:0 0 18px rgba(92,200,232,.25); }
.authwrap{ position:relative; z-index:1; }
</style>
"""

LOGIN_HERO_HTML = """
<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
@import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@500;700&family=JetBrains+Mono:wght@400;500&display=swap');
*{margin:0;padding:0;box-sizing:border-box;}
html,body{height:100%;overflow:hidden;background:transparent;}
#c{position:absolute;inset:0;cursor:crosshair;}
.sub{position:absolute;left:0;right:0;bottom:34px;text-align:center;
  font-family:'JetBrains Mono',monospace;font-size:12px;color:#8B96AB;letter-spacing:.22em;text-transform:uppercase;}
.sub .cur{color:#5CC8E8;}
.hint{position:absolute;left:0;right:0;bottom:12px;text-align:center;
  font-family:'JetBrains Mono',monospace;font-size:10px;color:#5A6679;letter-spacing:.14em;}
.corner{position:absolute;width:18px;height:18px;border:1px solid #2C3A52;opacity:.85;}
.corner.tl{top:10px;left:10px;border-right:none;border-bottom:none;}
.corner.tr{top:10px;right:10px;border-left:none;border-bottom:none;}
.corner.bl{bottom:10px;left:10px;border-right:none;border-top:none;}
.corner.br{bottom:10px;right:10px;border-left:none;border-top:none;}
</style></head><body>
<canvas id="c"></canvas>
<div class="corner tl"></div><div class="corner tr"></div><div class="corner bl"></div><div class="corner br"></div>
<div class="sub"><span id="t"></span><span class="cur">_</span></div>
<div class="hint">passa o rato pelas particulas &middot; inicia sessao abaixo</div>
<script>
var canvas=document.getElementById('c'),ctx=canvas.getContext('2d');
var W,H,P=[],mouse={x:-9999,y:-9999};
var COLS=['#5CC8E8','#5CC8E8','#E8EDF5','#D2A24C','#46B98C'];
function build(){
  W=canvas.width=canvas.offsetWidth; H=canvas.height=canvas.offsetHeight;
  var off=document.createElement('canvas'); off.width=W; off.height=H;
  var o=off.getContext('2d');
  var fs=Math.min(W/6.4,118);
  o.fillStyle='#fff';
  o.font='700 '+fs+'px "Space Grotesk", Arial, sans-serif';
  o.textAlign='center'; o.textBaseline='middle';
  o.fillText('TIC COPILOT', W/2, H/2-16);
  var img=o.getImageData(0,0,W,H).data;
  var targets=[], gap=Math.max(3,Math.round(fs/30));
  for(var y=0;y<H;y+=gap){ for(var x=0;x<W;x+=gap){
    if(img[(y*W+x)*4+3]>140) targets.push([x,y]);
  }}
  while(targets.length>3200){ targets.splice(Math.floor(Math.random()*targets.length),1); }
  P=targets.map(function(t,i){
    var edge=Math.random();
    var sx=edge<.5?Math.random()*W:(Math.random()<.5?-30:W+30);
    var sy=edge<.5?(Math.random()<.5?-30:H+30):Math.random()*H;
    return {x:sx,y:sy,vx:0,vy:0,tx:t[0],ty:t[1],c:COLS[i%COLS.length],s:1.3+Math.random()*1.1};
  });
}
build();
window.addEventListener('resize',build);
if(document.fonts&&document.fonts.ready){document.fonts.ready.then(build);}
canvas.addEventListener('mousemove',function(e){var r=canvas.getBoundingClientRect();mouse.x=e.clientX-r.left;mouse.y=e.clientY-r.top;});
canvas.addEventListener('mouseleave',function(){mouse.x=-9999;mouse.y=-9999;});
var RADIUS=90;
function step(){
  ctx.clearRect(0,0,W,H);
  for(var i=0;i<P.length;i++){
    var p=P[i];
    p.vx+=(p.tx-p.x)*0.022; p.vy+=(p.ty-p.y)*0.022;
    var dx=p.x-mouse.x, dy=p.y-mouse.y, d=Math.sqrt(dx*dx+dy*dy);
    if(d<RADIUS&&d>0.1){ var f=(RADIUS-d)/RADIUS; p.vx+=(dx/d)*f*3.4; p.vy+=(dy/d)*f*3.4; }
    p.vx*=0.86; p.vy*=0.86;
    p.x+=p.vx; p.y+=p.vy;
    ctx.globalAlpha=0.92;
    ctx.fillStyle=p.c;
    ctx.fillRect(p.x,p.y,p.s,p.s);
  }
  ctx.globalAlpha=1;
  requestAnimationFrame(step);
}
step();
var phrases=["Autoeuropa \u00b7 Technical Innovation Center","CAD \u00b7 3D Printing \u00b7 3D Reality",
  "Claude Opus 4.7 \u00d7 GPT-5.5"];
var pi=0,ci=0,del=false,el=document.getElementById('t');
function type(){var w=phrases[pi];el.textContent=w.substring(0,ci);
  if(!del){ci++;if(ci>w.length){del=true;setTimeout(type,1700);return;}}
  else{ci--;if(ci<0){del=false;pi=(pi+1)%phrases.length;ci=0;}}
  setTimeout(type,del?24:50);}
type();
</script></body></html>
"""


# ---------------------------------------------------------------------------
# Helpers de UI
# ---------------------------------------------------------------------------
def model_badge(k: str, sub: str | None = None) -> str:
    m = MODELS[k]
    sub_html = f"<span class='sub'>{sub}</span>" if sub else ""
    return f"<span class='badge badge--{k}'><span class='dot'></span>{m['label']}{sub_html}</span>"


def tic_badge(text: str) -> str:
    return f"<span class='badge badge--tic'><span class='dot'></span>{text}</span>"


def reason_card(k: str, text: str) -> str:
    return f"<div class='reason reason--{k}'>{text}</div>"


def metric_html(k: str, output: str, secs: float, prompt: str) -> str:
    return (f"<div class='metricline'><span>~{est_tokens(output)} tok</span>"
            f"<span class='sep'>/</span><span>{secs:.1f}s</span>"
            f"<span class='sep'>/</span><span>~${est_cost(prompt, output, k):.4f}</span></div>")


def step_html(n: int, text: str) -> str:
    return f"<div class='step'><span class='n'>{n}</span>{text}</div>"


def attach_note(names: list[str]) -> str:
    return f"<div class='attach-note'>anexos: {', '.join(names)}</div>"


def sources_note(sources: list[dict]) -> str:
    links = " · ".join(
        f"<a href='{s['url']}' target='_blank'>[{s['n']}] {s['title'][:48]}</a>"
        for s in sources if s.get("url"))
    return f"<div class='src-note'>fontes web: {links}</div>"


# ---------------------------------------------------------------------------
# Render de historico
# ---------------------------------------------------------------------------
def render_history():
    for msg in st.session_state.messages:
        if msg["role"] == "user":
            with st.chat_message("user"):
                st.markdown(msg["content"])
                if msg.get("attachments"):
                    st.markdown(attach_note(msg["attachments"]), unsafe_allow_html=True)
                if msg.get("web_sources"):
                    st.markdown(sources_note(msg["web_sources"]), unsafe_allow_html=True)
            continue
        with st.chat_message("assistant"):
            mode = msg.get("mode")
            if mode == "single":
                head = model_badge(msg["model"])
                if msg.get("reason"):
                    head += reason_card(msg["model"], msg["reason"])
                st.markdown(head, unsafe_allow_html=True)
                st.markdown(msg["content"])
            elif mode == "duel":
                cols = st.columns(2)
                for col, k in zip(cols, ["claude", "gpt"]):
                    with col:
                        st.markdown(f"<div class='col-head'>{model_badge(k)}</div>", unsafe_allow_html=True)
                        with st.container(border=True):
                            st.markdown(msg["data"].get(k, ""))
                        if msg["meta"].get(k):
                            st.markdown(msg["meta"][k], unsafe_allow_html=True)
            elif mode == "collab":
                ex, rv = msg["executor"], msg["reviewer"]
                st.markdown("Executor " + model_badge(ex) + " &nbsp; Revisor " + model_badge(rv),
                            unsafe_allow_html=True)
                st.markdown(step_html(1, f"Rascunho - {MODELS[ex]['label']}"), unsafe_allow_html=True)
                st.markdown(msg["draft"])
                st.markdown(step_html(2, f"Revisao critica - {MODELS[rv]['label']}"), unsafe_allow_html=True)
                st.markdown(msg["critique"])
                if msg.get("synthesis"):
                    st.markdown(step_html(3, f"Versao final - {MODELS[ex]['label']}"), unsafe_allow_html=True)
                    st.markdown(msg["synthesis"])


# ---------------------------------------------------------------------------
# Modos de resposta
# ---------------------------------------------------------------------------
def handle_single(send_text, display_text, forced_model, show_reason):
    if forced_model:
        model_key, reason = forced_model, "Modelo selecionado manualmente."
    else:
        decision = route(display_text, total_len=len(send_text))
        model_key, reason = decision["model"], decision["reason"]
    with st.chat_message("assistant"):
        head = model_badge(model_key)
        if show_reason:
            head += reason_card(model_key, reason)
        st.markdown(head, unsafe_allow_html=True)
        full = st.write_stream(smooth_stream(respond(
            model_key, send_text, st.session_state.threads[model_key],
            api_key_for(model_key), persist=True)))
    track_usage(model_key, send_text, full or "")
    st.session_state.messages.append({"role": "assistant", "mode": "single", "model": model_key,
                                      "reason": reason if show_reason else "", "content": full})


def handle_duel(send_text, display_text):
    order = ["claude", "gpt"]
    with st.chat_message("assistant"):
        cols = st.columns(2)
        ph, metaph = {}, {}
        for col, k in zip(cols, order):
            with col:
                st.markdown(f"<div class='col-head'>{model_badge(k)}</div>", unsafe_allow_html=True)
                with st.container(border=True):
                    ph[k] = st.empty()
                metaph[k] = st.empty()
        tids = {k: st.session_state.threads[k] for k in order}
        keys = {k: api_key_for(k) for k in order}
        qs = {k: queue.Queue() for k in order}
        results = {k: "" for k in order}
        starts = {k: time.time() for k in order}
        times = {k: 0.0 for k in order}

        def worker(k):
            try:
                for tok in smooth_stream(respond(k, send_text, tids[k], keys[k],
                                                 persist=False), delay=0.0):
                    qs[k].put(tok)
            except Exception as e:  # noqa
                qs[k].put(f"\n\n> Erro: {e}")
            finally:
                qs[k].put(None)

        threads = {k: threading.Thread(target=worker, args=(k,), daemon=True) for k in order}
        for t in threads.values():
            t.start()
        done = {k: False for k in order}
        while not all(done.values()):
            for k in order:
                updated = False
                try:
                    for _ in range(16):  # lote pequeno por frame -> fluxo suave
                        item = qs[k].get_nowait()
                        if item is None:
                            done[k] = True
                            times[k] = time.time() - starts[k]
                            break
                        results[k] += item
                        updated = True
                except queue.Empty:
                    pass
                if updated:
                    ph[k].markdown(results[k] + " \u258c")
            time.sleep(0.03)
        meta = {}
        for k in order:
            ph[k].markdown(results[k])
            meta[k] = metric_html(k, results[k], times[k], send_text)
            metaph[k].markdown(meta[k], unsafe_allow_html=True)
            track_usage(k, send_text, results[k])
    st.session_state.messages.append({"role": "assistant", "mode": "duel", "data": results, "meta": meta})


def handle_collab(send_text, display_text, do_synthesis):
    decision = route(display_text, total_len=len(send_text))
    executor = decision["model"]
    reviewer = "gpt" if executor == "claude" else "claude"
    ex_key, rv_key = api_key_for(executor), api_key_for(reviewer)
    with st.chat_message("assistant"):
        st.markdown("Executor " + model_badge(executor) + " &nbsp; Revisor " + model_badge(reviewer),
                    unsafe_allow_html=True)
        st.markdown(reason_card(executor, decision["reason"]), unsafe_allow_html=True)
        st.markdown(step_html(1, f"Rascunho - {MODELS[executor]['label']}"), unsafe_allow_html=True)
        draft = st.write_stream(smooth_stream(respond(executor, send_text, new_thread_id(), ex_key)))
        critique_prompt = (
            "Es um revisor tecnico critico e rigoroso. Outro assistente respondeu ao pedido "
            "abaixo. Identifica erros factuais, lacunas, ambiguidades, riscos e suposicoes "
            "frageis, e propoe melhorias CONCRETAS e acionaveis. Se direto; nao repitas tudo.\n\n"
            f"=== PEDIDO ===\n{display_text}\n\n=== RESPOSTA A REVER ===\n{draft}"
        )
        st.markdown(step_html(2, f"Revisao critica - {MODELS[reviewer]['label']}"), unsafe_allow_html=True)
        critique = st.write_stream(smooth_stream(respond(reviewer, critique_prompt, new_thread_id(), rv_key)))
        synthesis = ""
        if do_synthesis:
            synth_prompt = (
                "Tens o teu rascunho e uma revisao critica. Produz a VERSAO FINAL melhorada, "
                "integrando o feedback valido. Entrega so a resposta final.\n\n"
                f"=== PEDIDO ===\n{display_text}\n\n=== RASCUNHO ===\n{draft}\n\n"
                f"=== REVISAO ===\n{critique}"
            )
            st.markdown(step_html(3, f"Versao final - {MODELS[executor]['label']}"), unsafe_allow_html=True)
            synthesis = st.write_stream(smooth_stream(respond(executor, synth_prompt, new_thread_id(), ex_key)))
    track_usage(executor, send_text, (draft or "") + (synthesis or ""))
    track_usage(reviewer, critique_prompt, critique or "")
    st.session_state.messages.append({"role": "assistant", "mode": "collab", "executor": executor,
                                      "reviewer": reviewer, "draft": draft, "critique": critique,
                                      "synthesis": synthesis})


# ---------------------------------------------------------------------------
# Sessao / Auth UI
# ---------------------------------------------------------------------------
def reset_session():
    st.session_state.messages = []
    st.session_state.threads = {k: new_thread_id() for k in MODELS}
    st.session_state.uploader_key = st.session_state.get("uploader_key", 0) + 1
    st.session_state.audio_key = st.session_state.get("audio_key", 0) + 1
    st.session_state.conv_id = None
    st.session_state.conv_title = None
    st.session_state.conv_created = None
    st.session_state.pop("voice_draft", None)
    clear_session_rag()


def render_auth():
    users = _load_users()
    st.markdown("<div class='authwrap'>", unsafe_allow_html=True)
    if not users:
        st.markdown("<h3>Criar conta de administrador</h3>"
                    "<p class='s'>Primeiro arranque: define a tua conta.</p>", unsafe_allow_html=True)
        with st.form("setup"):
            name = st.text_input("Nome")
            username = st.text_input("Utilizador (a-z, 0-9, _ . -)")
            p1 = st.text_input("Password", type="password")
            p2 = st.text_input("Confirmar password", type="password")
            ok = st.form_submit_button("Criar conta", use_container_width=True)
        if ok:
            username = (username or "").strip().lower()
            if not USERNAME_RE.match(username):
                st.error("Utilizador invalido (3-32 chars: a-z, 0-9, _ . -).")
            elif len(p1) < 6:
                st.error("Password com pelo menos 6 caracteres.")
            elif p1 != p2:
                st.error("As passwords nao coincidem.")
            else:
                create_user(username, name or username, p1)
                st.session_state.auth = {"user": username, "name": name or username}
                reset_session()
                st.rerun()
    else:
        st.markdown("<h3>Entrar</h3><p class='s'>Acede com as tuas credenciais.</p>",
                    unsafe_allow_html=True)
        with st.form("login"):
            username = st.text_input("Utilizador")
            password = st.text_input("Password", type="password")
            ok = st.form_submit_button("Entrar", use_container_width=True)
        if ok:
            username = (username or "").strip().lower()
            if verify_user(username, password):
                st.session_state.auth = {"user": username, "name": _load_users()[username]["name"]}
                reset_session()
                st.rerun()
            else:
                st.error("Credenciais invalidas.")
        with st.expander("Criar nova conta"):
            with st.form("register"):
                name = st.text_input("Nome", key="r_name")
                u = st.text_input("Utilizador", key="r_user")
                p1 = st.text_input("Password", type="password", key="r_p1")
                p2 = st.text_input("Confirmar", type="password", key="r_p2")
                rok = st.form_submit_button("Registar", use_container_width=True)
            if rok:
                u = (u or "").strip().lower()
                if not USERNAME_RE.match(u):
                    st.error("Utilizador invalido.")
                elif u in users:
                    st.error("Esse utilizador ja existe.")
                elif len(p1) < 6 or p1 != p2:
                    st.error("Password fraca ou nao coincide.")
                else:
                    create_user(u, name or u, p1)
                    st.success("Conta criada. Ja podes entrar.")
    st.markdown("</div>", unsafe_allow_html=True)


# ---------------------------------------------------------------------------
# Sidebar
# ---------------------------------------------------------------------------
def sidebar(username):
    with st.sidebar:
        u = st.session_state.auth
        initial = (u["name"][:1] or "?").upper()
        st.markdown(f"<div class='userbox'><div class='av'>{initial}</div>"
                    f"<div><div class='n'>{u['name']}</div><div class='r'>@{u['user']}</div></div></div>",
                    unsafe_allow_html=True)
        if st.button("Terminar sessao", use_container_width=True):
            save_current_conversation(username)
            st.session_state.auth = None
            reset_session()
            st.rerun()

        # ---------------- Conversas (historico persistente) ----------------
        st.markdown("<div class='eyebrow'>Conversas</div>", unsafe_allow_html=True)
        if st.button("+ Nova conversa", use_container_width=True):
            save_current_conversation(username)
            reset_session()
            st.rerun()
        convs = list_conversations(username)
        active_id = st.session_state.get("conv_id")
        if not convs:
            st.caption("Sem conversas guardadas. A primeira fica gravada automaticamente.")
        for c in convs[:MAX_CONVERSATIONS_LISTED]:
            cols = st.columns([0.84, 0.16])
            mark = "» " if c["id"] == active_id else ""
            title = c["title"] if len(c["title"]) <= 34 else c["title"][:34] + "..."
            if cols[0].button(f"{mark}{title}", key=f"conv_{c['id']}", use_container_width=True,
                              help=f"{c['n']} mensagens · {c['updated'][:16]}"):
                save_current_conversation(username)
                if load_conversation(username, c["id"]):
                    st.rerun()
            if cols[1].button("✕", key=f"del_{c['id']}", help="Apagar conversa"):
                delete_conversation(username, c["id"])
                if c["id"] == active_id:
                    reset_session()
                st.rerun()
        if st.session_state.get("messages"):
            st.download_button("Exportar conversa (.md)", conversation_markdown(),
                               file_name=f"conversa_{st.session_state.get('conv_id') or 'atual'}.md",
                               mime="text/markdown", use_container_width=True)

        # ---------------- Modo ----------------
        st.markdown("<div class='eyebrow'>Modo de operacao</div>", unsafe_allow_html=True)
        mode = st.radio("Modo", ["Auto (Roteador)", "So Claude", "So GPT",
                                 "Duelo (lado a lado)", "Colaboracao"],
                        index=0, label_visibility="collapsed")
        show_reason = st.toggle("Mostrar justificacao do roteamento", value=True)
        do_synthesis = st.toggle("Sintese final (colaboracao)", value=True) if mode == "Colaboracao" else False

        # ---------------- Pesquisa web ----------------
        st.markdown("<div class='eyebrow'>Pesquisa web</div>", unsafe_allow_html=True)
        if HAS_DDG:
            web_mode = st.radio("Web", ["Auto (roteador)", "Sempre", "Nunca"],
                                index=0, label_visibility="collapsed", horizontal=True)
        else:
            web_mode = "Nunca"
            st.caption("Pesquisa web off: instala `ddgs` (pip install ddgs).")

        # ---------------- Anexos ----------------
        st.markdown("<div class='eyebrow'>Anexos (PDF, PPTX, imagem)</div>", unsafe_allow_html=True)
        files = st.file_uploader("Anexar", type=["pdf", "pptx", "ppt", "png", "jpg", "jpeg", "webp", "bmp", "tiff", "txt", "md", "csv"],
                                 accept_multiple_files=True, label_visibility="collapsed",
                                 key=f"uploader_{st.session_state.uploader_key}")
        if files:
            st.caption(f"{len(files)} ficheiro(s) sera(o) anexado(s) ao proximo pedido.")
        use_vlm = st.toggle("Interpretar imagens com VLM",
                            value=bool(hf_token()),
                            help=f"Modelo gratuito ({HF_VLM_MODEL}) via Hugging Face. Requer HF_TOKEN.")
        if use_vlm and not hf_token():
            st.caption("Define `HF_TOKEN` (gratuito em huggingface.co/settings/tokens) para ativar o VLM.")
        if not HAS_OCR:
            st.caption("OCR off: instala 'tesseract-ocr' p/ ler texto em imagens e PDFs digitalizados.")

        # ---------------- RAG ----------------
        st.markdown("<div class='eyebrow'>RAG hibrido dos anexos</div>", unsafe_allow_html=True)
        use_attachment_search = st.toggle("Usar RAG hibrido (BM25 + embeddings + RRF)", value=True)
        attachment_top_k = st.slider("Trechos RAG a enviar", 2, 16, ATTACHMENT_TOP_K_DEFAULT, 1)
        embedding_model_name = st.selectbox("Modelo de embeddings", EMBEDDING_MODEL_OPTIONS, index=0)
        use_reranker = st.toggle("Reranker neural", value=False,
                                 help="Cross-encoder: mais preciso em documentos longos, mas mais lento.")
        reranker_model_name = st.selectbox("Modelo reranker", RERANKER_MODEL_OPTIONS, index=0,
                                           disabled=not use_reranker)
        _ensure_session_rag_docs()
        if st.session_state.rag_docs:
            total_chars = sum(len(d.get("text", "")) for d in st.session_state.rag_docs)
            st.caption(f"Indice da conversa: {len(st.session_state.rag_docs)} doc(s), "
                       f"{total_chars:,} caracteres".replace(",", " "))
            if st.button("Limpar indice RAG", use_container_width=True):
                clear_session_rag()
                st.rerun()
        missing = []
        if not HAS_SENTENCE_TRANSFORMERS:
            missing.append("sentence-transformers")
        if not HAS_FAISS:
            missing.append("faiss-cpu")
        if not HAS_NUMPY:
            missing.append("numpy")
        if not HAS_BM25:
            missing.append("rank-bm25")
        if missing:
            st.caption("RAG parcial/indisponivel sem: " + ", ".join(missing))

        # ---------------- Voz ----------------
        st.markdown("<div class='eyebrow'>Voz</div>", unsafe_allow_html=True)
        if HAS_WHISPER:
            whisper_size = st.selectbox("Modelo de transcricao (Whisper)", WHISPER_SIZES,
                                        index=WHISPER_SIZES.index(WHISPER_DEFAULT))
        else:
            whisper_size = WHISPER_DEFAULT
            st.caption("Voz off: instala `faster-whisper` (pip install faster-whisper).")

        # ---------------- Memoria ----------------
        st.markdown("<div class='eyebrow'>Memoria</div>", unsafe_allow_html=True)
        use_memory = st.toggle("Usar memoria master no contexto", value=True)
        use_session_context = st.toggle("Usar conversa atual como contexto", value=True)
        auto_mem = st.toggle("Atualizar memoria automaticamente", value=True)
        mem = load_memory(username)
        with st.expander("Ver / editar memoria master"):
            edited = st.text_area("Resumo", value=mem.get("master", ""), height=160,
                                  label_visibility="collapsed")
            c1, c2 = st.columns(2)
            if c1.button("Guardar", use_container_width=True):
                mem["master"] = edited
                mem["updated_at"] = dt.datetime.now().isoformat(timespec="seconds")
                save_memory(username, mem)
                st.toast("Memoria guardada.")
            if c2.button("Atualizar agora", use_container_width=True):
                with st.spinner("A resumir..."):
                    summarize_memory(username)
                st.rerun()
            if mem.get("updated_at"):
                st.caption(f"Atualizada: {mem['updated_at']}")
            if st.button("Limpar memoria", use_container_width=True):
                save_memory(username, {"master": "", "updated_at": None, "turns_since": 0})
                st.rerun()

        # ---------------- Sessao / custos ----------------
        st.markdown("<div class='eyebrow'>Sessao</div>", unsafe_allow_html=True)
        usage = st.session_state.get("usage", {"cost": 0.0, "calls": 0, "tokens_out": 0})
        st.markdown(f"<div class='usagebox'><span>custo <b>~${usage['cost']:.4f}</b></span>"
                    f"<span>{usage['calls']} resp.</span>"
                    f"<span>~{usage['tokens_out']} tok out</span></div>", unsafe_allow_html=True)

        st.markdown("<div class='eyebrow'>Avancado</div>", unsafe_allow_html=True)
        _min_interval[0] = st.slider("Intervalo minimo entre pedidos (s)", 0.0, 5.0, _min_interval[0], 0.1)

        with st.expander("Como funciona o roteamento"):
            rows = ""
            for r in ROUTING_RULES:
                c = MODELS[r["model"]]["color"]
                rows += (f"<div class='routerule'><span class='d' style='background:{c}'></span>"
                         f"<div>{r['label']}<span class='b'>{MODELS[r['model']]['label']} &middot; {r['bench']}</span></div></div>")
            c = MODELS[DEFAULT_MODEL]["color"]
            rows += (f"<div class='routerule'><span class='d' style='background:{c}'></span>"
                     f"<div>Geral / sem sinal<span class='b'>{MODELS[DEFAULT_MODEL]['label']} &middot; default</span></div></div>")
            st.markdown(rows, unsafe_allow_html=True)

    return (mode, show_reason, do_synthesis, files, use_memory, use_session_context, auto_mem,
            use_attachment_search, attachment_top_k, embedding_model_name, use_reranker,
            reranker_model_name, web_mode, use_vlm, whisper_size)


# ---------------------------------------------------------------------------
# Barra de entrada: acoes rapidas + voz
# ---------------------------------------------------------------------------
def render_quick_actions():
    """Acoes rapidas TIC — so quando a conversa esta vazia."""
    if st.session_state.get("messages"):
        return
    st.markdown("<div class='eyebrow'>Comecar com uma acao do TIC</div>", unsafe_allow_html=True)
    cols = st.columns(2)
    for i, (label, prompt) in enumerate(QUICK_ACTIONS):
        if cols[i % 2].button(label, key=f"qa_{i}", use_container_width=True):
            st.session_state.pending_prompt = prompt
            st.rerun()


def render_voice_input(whisper_size: str):
    """Gravacao no browser + transcricao local com faster-whisper."""
    with st.popover("🎙️ Entrada por voz", use_container_width=False):
        if not HAS_WHISPER:
            st.caption("Instala `faster-whisper` para transcrever audio localmente (gratis).")
            return
        audio = st.audio_input("Grava o teu pedido",
                               key=f"audio_{st.session_state.get('audio_key', 0)}")
        if audio is not None:
            sig = hashlib.sha1(audio.getvalue()).hexdigest()
            if st.session_state.get("voice_sig") != sig:
                with st.spinner("A transcrever (Whisper local)..."):
                    try:
                        st.session_state.voice_draft = transcribe_audio(audio.getvalue(), whisper_size)
                        st.session_state.voice_sig = sig
                    except Exception as e:
                        st.error(f"Falha na transcricao: {e}")
        draft = st.session_state.get("voice_draft", "")
        if draft:
            edited = st.text_area("Transcricao (edita se precisares)", value=draft, height=110)
            c1, c2 = st.columns(2)
            if c1.button("Enviar pedido", use_container_width=True, type="primary"):
                st.session_state.pending_prompt = edited.strip()
                st.session_state.pop("voice_draft", None)
                st.session_state.pop("voice_sig", None)
                st.session_state.audio_key = st.session_state.get("audio_key", 0) + 1
                st.rerun()
            if c2.button("Descartar", use_container_width=True):
                st.session_state.pop("voice_draft", None)
                st.session_state.pop("voice_sig", None)
                st.session_state.audio_key = st.session_state.get("audio_key", 0) + 1
                st.rerun()


# ---------------------------------------------------------------------------
# UI principal
# ---------------------------------------------------------------------------
def main():
    st.set_page_config(page_title="TIC Copilot — Claude x GPT", page_icon="◆", layout="wide")
    if "uploader_key" not in st.session_state:
        st.session_state.uploader_key = 0
    if "audio_key" not in st.session_state:
        st.session_state.audio_key = 0
    if "auth" not in st.session_state:
        st.session_state.auth = None
    st.markdown(THEME_CSS, unsafe_allow_html=True)

    if not st.session_state.auth:
        # Pagina de login: aurora animada + particulas interativas "TIC COPILOT"
        st.markdown(LOGIN_CSS, unsafe_allow_html=True)
        components.html(LOGIN_HERO_HTML, height=330, scrolling=False)
        render_auth()
        return

    # Hero personalizado: "Bom dia / Boa tarde / Boa noite, <nome>"
    greet = greeting_pt(st.session_state.auth.get("name", ""))
    components.html(hero_html(greet, date_pt()), height=240, scrolling=False)

    username = st.session_state.auth["user"]
    if "messages" not in st.session_state:
        reset_session()

    (mode, show_reason, do_synthesis, files, use_memory, use_session_context, auto_mem,
     use_attachment_search, attachment_top_k, embedding_model_name, use_reranker,
     reranker_model_name, web_mode, use_vlm, whisper_size) = sidebar(username)

    render_history()
    render_quick_actions()
    render_voice_input(whisper_size)

    prompt = st.chat_input("Escreve o teu pedido... (ou usa a entrada por voz)")
    pending = st.session_state.pop("pending_prompt", None)
    if not prompt and pending:
        prompt = pending
    if not prompt:
        return

    # ---------------- Anexos (extracao + RAG hibrido) ----------------
    attach_text, attach_names, attach_meta = "", [], {"mode": "empty"}
    if files:
        with st.spinner("A processar anexos (extracao / OCR / VLM)..."):
            attach_text, attach_names, warns, attach_meta = extract_files(
                files,
                query=prompt,
                use_attachment_search=use_attachment_search,
                attachment_top_k=attachment_top_k,
                embedding_model_name=embedding_model_name,
                use_reranker=use_reranker,
                reranker_model_name=reranker_model_name,
                use_vlm=use_vlm,
            )
        for w in warns:
            st.warning(w)
        if attach_meta.get("mode") == "hybrid_rag":
            vias = []
            if attach_meta.get("dense"):
                vias.append("densa")
            if attach_meta.get("sparse"):
                vias.append("BM25")
            st.info(f"RAG hibrido ({' + '.join(vias)} + RRF): "
                    f"{attach_meta.get('total_chunks', 0)} chunks avaliados; "
                    f"{len(attach_meta.get('selected', []))} enviados ao modelo.")
    elif use_attachment_search and st.session_state.get("rag_docs"):
        with st.spinner("A recuperar trechos relevantes do indice RAG da conversa..."):
            try:
                attach_text, attach_meta = retrieve_session_rag_context(
                    prompt,
                    attachment_top_k=attachment_top_k,
                    embedding_model_name=embedding_model_name,
                    use_reranker=use_reranker,
                    reranker_model_name=reranker_model_name,
                )
                attach_names = [d.get("name", "documento") for d in st.session_state.rag_docs]
            except Exception as e:
                st.warning(f"Falha no RAG da conversa: {e}")
                attach_text, attach_meta = "", {"mode": "error", "error": str(e)}

    # ---------------- Pesquisa web ----------------
    web_text, web_sources = "", []
    decision_preview = route(prompt, total_len=len(prompt))
    should_search = (web_mode == "Sempre"
                     or (web_mode == "Auto (roteador)" and decision_preview.get("rule_id") == "web"))
    if should_search and HAS_DDG:
        with st.spinner("A pesquisar na web (DuckDuckGo)..."):
            web_text, web_sources = web_search_context(prompt)
        if not web_text:
            st.warning("Pesquisa web sem resultados utilizaveis - a responder so com o modelo.")

    display_text = prompt
    send_text = build_agent_message(prompt, username, use_memory, attach_text,
                                    use_session_context, web_text=web_text)

    st.session_state.messages.append({"role": "user", "content": prompt,
                                      "attachments": attach_names,
                                      "web_sources": web_sources})
    with st.chat_message("user"):
        st.markdown(prompt)
        if attach_names:
            st.markdown(attach_note(attach_names), unsafe_allow_html=True)
        if web_sources:
            st.markdown(sources_note(web_sources), unsafe_allow_html=True)

    if mode == "Auto (Roteador)":
        handle_single(send_text, display_text, None, show_reason)
    elif mode == "So Claude":
        handle_single(send_text, display_text, "claude", show_reason)
    elif mode == "So GPT":
        handle_single(send_text, display_text, "gpt", show_reason)
    elif mode == "Duelo (lado a lado)":
        handle_duel(send_text, display_text)
    elif mode == "Colaboracao":
        handle_collab(send_text, display_text, do_synthesis)

    # Historico persistente: grava a conversa apos cada turno
    save_current_conversation(username)

    # Memoria: log + atualizacao automatica
    append_log(username, display_text, mode)
    mem = load_memory(username)
    mem["turns_since"] = mem.get("turns_since", 0) + 1
    save_memory(username, mem)
    if auto_mem and mem["turns_since"] >= MEMORY_UPDATE_EVERY:
        with st.spinner("A atualizar memoria..."):
            summarize_memory(username)

    # Limpa o uploader para o proximo turno
    if attach_names and files:
        st.session_state.uploader_key += 1
        st.rerun()


if __name__ == "__main__":
    main()
